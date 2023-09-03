
import pandas as pd
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

from pptx.compat import BytesIO
from pptx.enum.shapes import PP_PLACEHOLDER, PROG_ID
from pptx.media import SPEAKER_IMAGE_BYTES, Video
from pptx.opc.constants import CONTENT_TYPE as CT
from pptx.oxml.ns import qn
from pptx.oxml.shapes.graphfrm import CT_GraphicalObjectFrame
from pptx.oxml.shapes.picture import CT_Picture
from pptx.oxml.simpletypes import ST_Direction
from pptx.shapes.autoshape import AutoShapeType, Shape
from pptx.shapes.base import BaseShape
from pptx.shapes.connector import Connector
from pptx.shapes.freeform import FreeformBuilder
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.group import GroupShape
from pptx.shapes.picture import Movie, Picture
from pptx.enum.text import PP_ALIGN

powerpoint_file_path = "/Users/erick/Desktop/Graphics_Templates/email_template.pptx"



csv_file_path = "/Users/erick/Desktop/day_part_data.csv"
csv_file_path2 = "/Users/erick/Desktop/city_high_and Lows.csv"
csv_file_path3 = "/Users/erick/Desktop/new york_7_day_forecast.csv"


# Set the slide index and text box index of the PowerPoint slide to update
 # REMINDER: Slide index is 0-based, so slide 7 corresponds to index 6
slide_index = 0
slide_index2 = 1
slide_index3 = 2
slide_index4 = 3
slide_index5 = 4

#READ THE DATA
data = pd.read_csv(csv_file_path)
data2 = pd.read_csv(csv_file_path2)
data3 = pd.read_csv(csv_file_path3)



#DAYPART TEMPS
daypart1_value = str(data.iloc[1, 2])
daypart2_value = str(data.iloc[1, 3])
daypart3_value = str(data.iloc[1, 4])
daypart4_value = str(data.iloc[1, 5])
daypart5_value = str(data.iloc[1, 6])

daypart6_value = str(data.iloc[3, 2])
daypart7_value = str(data.iloc[3, 3])
daypart8_value = str(data.iloc[3, 4])
daypart9_value = str(data.iloc[3, 5])
daypart10_value = str(data.iloc[3,6])

daypart11_value= str(data3.iloc[0, 2])
daypart12_value = str(data2.iloc[27,3])

daypart13_value = str(data2.iloc[22,2])
daypart14_value = str(data2.iloc[22,3])


#DAYPART WEATHER
daypart1_weather = str(data.iloc[0,1])
daypart2_weather = str(data.iloc[0,2])
daypart3_weather = str(data.iloc[0,3])
daypart4_weather = str(data.iloc[0,4])
daypart5_weather = str(data.iloc[0,5])

daypart6_weather = str(data.iloc[2,1])
daypart7_weather = str(data.iloc[2,2])
daypart8_weather = str(data.iloc[2,3])
daypart9_weather = str(data.iloc[2,4])
daypart10_weather = str(data.iloc[2,5])

daya_weather = str(data2.iloc[27,4])
dayb_weather = str(data2.iloc[22,4])

#NYC METRO HIGH TEMPS
cell_value_bx = str(data2.iloc[31, 2])
cell_value_jc = str(data2.iloc[19, 2])
cell_value_nk = str(data2.iloc[28, 2])
cell_value_si = str(data2.iloc[8, 2])
cell_value_bk = str(data2.iloc[6, 2])
cell_value_qn = str(data2.iloc[34, 2])
cell_value_gc = str(data2.iloc[17, 2])
cell_value_lb = str(data2.iloc[21, 2])
cell_value_bs = str(data2.iloc[3, 2])

#NYC METRO WEATHER
weather_bx = str(data2.iloc[31, 4])
weather_jc = str(data2.iloc[19, 4])
weather_nk = str(data2.iloc[28, 4])
weather_si = str(data2.iloc[8, 4])
weather_bk = str(data2.iloc[6, 4])
weather_qn = str(data2.iloc[34, 4])
weather_gc = str(data2.iloc[17, 4])
weather_lb = str(data2.iloc[21, 4])
weather_bs = str(data2.iloc[3, 4])

#NYC METRO RAIN TOTALS
rain_bx = str(data2.iloc[31, 5])
rain_jc = str(data2.iloc[19, 5])
rain_nk = str(data2.iloc[28, 5])
rain_si = str(data2.iloc[8, 5])
rain_bk = str(data2.iloc[6, 5])
rain_qn = str(data2.iloc[34, 5])
rain_gc = str(data2.iloc[17, 5])
rain_lb = str(data2.iloc[21, 5])
rain_bs = str(data2.iloc[3, 5])

#NYC METRO HEAT INDEX
app_value_bx = str(data2.iloc[31, 7])
app_value_jc = str(data2.iloc[19, 7])
app_value_nk = str(data2.iloc[28, 7])
app_value_si = str(data2.iloc[8, 7])
app_value_bk = str(data2.iloc[6, 7])
app_value_qn = str(data2.iloc[34, 7])
app_value_gc = str(data2.iloc[17, 7])
app_value_lb = str(data2.iloc[21, 7])
app_value_bs = str(data2.iloc[3, 7])

#IMAGE MAPPING
# Define the base directory for the image files
base_directory = "/Users/erick/Desktop/Weather_Graphics/Simple Weather Icons/weather_icons/"
base_directory2 = "/Users/erick/Desktop/Weather_Graphics/Simple Weather Icons/weather_icons_night/"

weather_image_mapping2 = {
    "sky is clear": "Moon + Stars.png",
    "moderate rain": "Rain.png",
    "light rain": "Rain.png",
    "overcast clouds": "Cloud.png",
    "scattered clouds": "Night + Clouds.png",
    "broken clouds": "Night + Clouds.png",
    "few clouds": "Moon + Stars.png",
    "heavy intensity rain": "Thunderstorm 2.png",
    "clear sky": "Moon + Stars.png",
    "partly cloudy": "Night + Clouds.png",
    "sunny": "Moon + Stars.png",
    "patchy rain possible": "Rain.png",
    "heavy rain": "Thunderstorm 2.png",
    "thunderstorm with rain": "Thunderstorm 2.png",
    "thunderstorm with heavy rain": "Thunderstorm 2.png"
}

# Define the dictionary mapping weather values to image file paths
weather_image_mapping = {
    "sky is clear": "Sun 3.png",
    "moderate rain": "Rain.png",
    "light rain": "Rain + Sun.png",
    "overcast clouds": "Cloud.png",
    "scattered clouds": "Sun & Clouds.png",
    "broken clouds": "Sun & Clouds.png",
    "few clouds": "Sun 3.png",
    "heavy intensity rain": "Thunderstorm & Sun.png",
    "clear sky": "Sun 3.png",
    "partly cloudy": "Sun & Clouds.png",
    "sunny": "Sun 3.png",
    "patchy rain possible": "Rain + Sun.png",
    "heavy rain": "Thunderstorm & Sun.png",
    "thunderstorm with rain": "Thunderstorm & Sun.png",
    "thunderstorm with heavy rain": "Thunderstorm 2.png"

    # Add more mappings for other weather conditions
}

daypart1_image_file = base_directory + weather_image_mapping.get(daypart1_weather, "Wind.png")
daypart2_image_file = base_directory + weather_image_mapping.get(daypart2_weather, "Wind.png")
daypart3_image_file = base_directory + weather_image_mapping.get(daypart3_weather, "Wind.png")
daypart4_image_file = base_directory + weather_image_mapping.get(daypart4_weather, "Wind.png")
daypart5_image_file = base_directory2 + weather_image_mapping2.get(daypart5_weather, "Wind.png")
daypart6_image_file = base_directory + weather_image_mapping.get(daypart6_weather, "Wind.png")
daypart7_image_file = base_directory + weather_image_mapping.get(daypart7_weather, "Wind.png")
daypart8_image_file = base_directory + weather_image_mapping.get(daypart8_weather, "Wind.png")
daypart9_image_file = base_directory + weather_image_mapping.get(daypart9_weather, "Wind.png")
daypart10_image_file = base_directory2 + weather_image_mapping2.get(daypart10_weather, "Wind.png")
day_image_file = base_directory + weather_image_mapping.get(daya_weather, "Wind.png")
dayb_image_file = base_directory + weather_image_mapping.get(dayb_weather, "Wind.png")





# SLIDE NUMBER 1 BEGINS HERE **************************
presentation = Presentation(powerpoint_file_path)

slide = presentation.slides[slide_index]

daypart1 = 4  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textboxa = slide.shapes[daypart1].text_frame

daypart2 = 7  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textboxb = slide.shapes[daypart2].text_frame

daypart3 = 10  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textboxc = slide.shapes[daypart3].text_frame

daypart4 = 13  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textboxd = slide.shapes[daypart4].text_frame

daypart5 = 16  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textboxe = slide.shapes[daypart5].text_frame

daypart11 = 21  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textboxf = slide.shapes[daypart11].text_frame

daypart12 = 22  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textboxg = slide.shapes[daypart12].text_frame

#CLEAR TEXT, ADD NEW VALUES
textboxa.clear()
textboxa.text = daypart1_value
textboxb.clear()
textboxb.text = daypart2_value
textboxc.clear()
textboxc.text = daypart3_value
textboxd.clear()
textboxd.text = daypart4_value
textboxe.clear()
textboxe.text = daypart5_value
textboxf.clear()
textboxf.text = daypart11_value
textboxg.clear()
textboxg.text = daypart12_value


#FORMATTING NEW TEXT


for paragraph in textboxa.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = False  
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textboxb.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = False  
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textboxc.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = False  
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textboxd.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = False  
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textboxe.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = False  
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textboxf.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(54)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = False  
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textboxg.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(54)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = False  
    paragraph.alignment = PP_ALIGN.CENTER

#presentation = Presentation(powerpoint_file_path)

for slide in presentation.slides:
    for shape in slide.shapes:
        if shape.name == "daypart1_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(daypart1_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "daypart2_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(daypart2_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "daypart3_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(daypart3_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "daypart4_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(daypart4_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "daypart5_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(daypart5_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "day_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(day_image_file, shape.left, shape.top, shape.width, shape.height)

# SLIDE NUMBER TWO BEGINS HERE

slide = presentation.slides[slide_index2]

daypart6 = 4  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textboxa = slide.shapes[daypart6].text_frame

daypart7 = 7  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textboxb = slide.shapes[daypart7].text_frame

daypart8 = 10  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textboxc = slide.shapes[daypart8].text_frame

daypart9 = 13  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textboxd = slide.shapes[daypart9].text_frame

daypart10 = 16  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textboxe = slide.shapes[daypart10].text_frame

daypart13 = 21  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textboxf = slide.shapes[daypart13].text_frame

daypart14 = 22  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textboxg = slide.shapes[daypart14].text_frame

#CLEAR TEXT, ADD NEW VALUES
textboxa.clear()
textboxa.text = daypart6_value
textboxb.clear()
textboxb.text = daypart7_value
textboxc.clear()
textboxc.text = daypart8_value
textboxd.clear()
textboxd.text = daypart9_value
textboxe.clear()
textboxe.text = daypart10_value
textboxf.clear()
textboxf.text = daypart13_value
textboxg.clear()
textboxg.text = daypart14_value


#FORMATTING NEW TEXT


for paragraph in textboxa.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = False  
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textboxb.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = False  
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textboxc.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = False  
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textboxd.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = False  
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textboxe.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = False  
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textboxf.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(54)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = False  
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textboxg.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(54)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = False  
    paragraph.alignment = PP_ALIGN.CENTER

#presentation = Presentation(powerpoint_file_path)

for slide in presentation.slides:
    for shape in slide.shapes:
        if shape.name == "daypart6_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(daypart6_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "daypart7_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(daypart7_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "daypart8_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(daypart8_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "daypart9_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(daypart9_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "daypart10_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(daypart10_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "dayb_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(dayb_image_file, shape.left, shape.top, shape.width, shape.height)

# SLIDE NUMBER 3 BEGINS HERE **************************


slide = presentation.slides[slide_index3]


bronx = 23  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox = slide.shapes[bronx].text_frame

jerseycity = 11  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox2 = slide.shapes[jerseycity].text_frame

newark = 25  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox3 = slide.shapes[newark].text_frame

statenisland = 17  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox4 = slide.shapes[statenisland].text_frame

brooklyn = 20  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox5 = slide.shapes[brooklyn].text_frame

queens = 8  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox6 = slide.shapes[queens].text_frame

gardencity = 24  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox7 = slide.shapes[gardencity].text_frame

longbeach = 14  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox8 = slide.shapes[longbeach].text_frame

bayshore = 36  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox9 = slide.shapes[bayshore].text_frame

# In[253]:

#CLEAR TEXT, ADD NEW VALUES
textbox.clear()
textbox.text = cell_value_bx
textbox2.clear()
textbox2.text = cell_value_jc
textbox3.clear()
textbox3.text = cell_value_nk
textbox4.clear()
textbox4.text = cell_value_si
textbox5.clear()
textbox5.text = cell_value_bk
textbox6.clear()
textbox6.text = cell_value_qn
textbox7.clear()
textbox7.text = cell_value_gc
textbox8.clear()
textbox8.text = cell_value_lb
textbox9.clear()
textbox9.text = cell_value_bs


#FORMATTING NEW TEXT


for paragraph in textbox.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox2.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox3.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox4.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox5.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER


for paragraph in textbox6.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox7.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox8.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox9.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

# Assign image file paths based on weather values using the dictionary mapping
bx_image_file = base_directory + weather_image_mapping.get(weather_bx, "Wind.png")
jc_image_file = base_directory + weather_image_mapping.get(weather_jc, "Wind.png")
nk_image_file = base_directory + weather_image_mapping.get(weather_nk, "Wind.png")
si_image_file = base_directory + weather_image_mapping.get(weather_si, "Wind.png")
bk_image_file = base_directory + weather_image_mapping.get(weather_bk, "Wind.png")
qn_image_file = base_directory + weather_image_mapping.get(weather_qn, "Wind.png")
gc_image_file = base_directory + weather_image_mapping.get(weather_gc, "Wind.png")
lb_image_file = base_directory + weather_image_mapping.get(weather_lb, "Wind.png")
bs_image_file = base_directory + weather_image_mapping.get(weather_bs, "Wind.png")


# Iterate through the slides and replace images
for slide in presentation.slides:
    for shape in slide.shapes:
        if shape.name == "bx_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(bx_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "jc_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(jc_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "nk_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(nk_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "si_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(si_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "bk_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(bk_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "qn_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(qn_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "gc_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(gc_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "lb_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(lb_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "bs_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(bs_image_file, shape.left, shape.top, shape.width, shape.height)



# SLIDE NUMBER 4 BEGINS HERE


slide = presentation.slides[slide_index4]


bronx = 23  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox = slide.shapes[bronx].text_frame

jerseycity = 22  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox2 = slide.shapes[jerseycity].text_frame

newark = 21  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox3 = slide.shapes[newark].text_frame

statenisland = 20  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox4 = slide.shapes[statenisland].text_frame

brooklyn = 28  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox5 = slide.shapes[brooklyn].text_frame

queens = 24  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox6 = slide.shapes[queens].text_frame

gardencity = 27  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox7 = slide.shapes[gardencity].text_frame

longbeach = 25  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox8 = slide.shapes[longbeach].text_frame

bayshore = 26  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox9 = slide.shapes[bayshore].text_frame

# In[253]:

#CLEAR TEXT, ADD NEW VALUES
textbox.clear()
textbox.text = f'{rain_bx}"'
textbox2.clear()
textbox2.text = f'{rain_jc}"'
textbox3.clear()
textbox3.text = f'{rain_nk}"'
textbox4.clear()
textbox4.text = f'{rain_si}"'
textbox5.clear()
textbox5.text = f'{rain_bk}"'
textbox6.clear()
textbox6.text = f'{rain_qn}"'
textbox7.clear()
textbox7.text = f'{rain_gc}"'
textbox8.clear()
textbox8.text = f'{rain_lb}"'
textbox9.clear()
textbox9.text = f'{rain_bs}"'


#FORMATTING NEW TEXT


for paragraph in textbox.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox2.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox3.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox4.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox5.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER


for paragraph in textbox6.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox7.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox8.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox9.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER



slide = presentation.slides[slide_index5]


bronx = 23  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox = slide.shapes[bronx].text_frame

jerseycity = 11  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox2 = slide.shapes[jerseycity].text_frame

newark = 25  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox3 = slide.shapes[newark].text_frame

statenisland = 17  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox4 = slide.shapes[statenisland].text_frame

brooklyn = 20  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox5 = slide.shapes[brooklyn].text_frame

queens = 8  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox6 = slide.shapes[queens].text_frame

gardencity = 24  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox7 = slide.shapes[gardencity].text_frame

longbeach = 14  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox8 = slide.shapes[longbeach].text_frame

bayshore = 28  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox9 = slide.shapes[bayshore].text_frame

# In[253]:

#CLEAR TEXT, ADD NEW VALUES
textbox.clear()
textbox.text = app_value_bx
textbox2.clear()
textbox2.text = app_value_jc
textbox3.clear()
textbox3.text = app_value_nk
textbox4.clear()
textbox4.text = app_value_si
textbox5.clear()
textbox5.text = app_value_bk
textbox6.clear()
textbox6.text = app_value_qn
textbox7.clear()
textbox7.text = app_value_gc
textbox8.clear()
textbox8.text = app_value_lb
textbox9.clear()
textbox9.text = app_value_bs


#FORMATTING NEW TEXT


for paragraph in textbox.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(139, 0, 0)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox2.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(139, 0, 0)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox3.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(139, 0, 0)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox4.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(139, 0, 0)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox5.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(139, 0, 0)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER


for paragraph in textbox6.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(139, 0, 0)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox7.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(139, 0, 0)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox8.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(139, 0, 0)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox9.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(139, 0, 0)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

#THIS UPDATES THE PRESENTATION
updated_powerpoint_file_path = "/Users/erick/Desktop/Wx Email Graphics.pptx"
presentation.save(updated_powerpoint_file_path)