
import pandas as pd
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

from pptx.compat import BytesIO
from pptx.enum.shapes import PP_PLACEHOLDER, PROG_ID
from pptx.media import SPEAKER_IMAGE_BYTES, Video
from pptx.opc.constants import CONTENT_TYPE as CT
from pptx.oxml.ns import qn
from pptx.oxml.shapes.graphfrm import CT_GraphicalObjectFrame
from pptx.oxml.shapes.picture import CT_Picture
from pptx.oxml.simpletypes import ST_Direction
from pptx.shapes.autoshape import AutoShapeType, Shape
from pptx.shapes.base import BaseShape
from pptx.shapes.connector import Connector
from pptx.shapes.freeform import FreeformBuilder
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.group import GroupShape
from pptx.shapes.picture import Movie, Picture
from pptx.enum.text import PP_ALIGN


powerpoint_file_path = "/Users/erick/Desktop/Graphics_Templates/weekend_weather_template.pptx"

csv_file_path = "/Users/erick/Desktop/city_high_and Lows.csv"
csv_file_path2 = "/Users/erick/Desktop/city_high_and Lows_added.csv"

#DATA SOURCE
data = pd.read_csv(csv_file_path)
data2 = pd.read_csv(csv_file_path2)

# Set the slide index and text box index of the PowerPoint slide to update

slide_index = 0 #NORTHEAST SATURDAY
slide_index2= 1 #NORTHEAST SUNDAY
slide_index3= 2 #WEST SATURDAY
slide_index4 = 3 #WEST SUNDAY
slide_index5 = 4 #MIDWEST SATURDAY
slide_index6 = 5 #MIDWEST SUNDAY
slide_index7 = 6 #SOUTH SATURDAY
slide_index8 = 7 #SOUTH SUNDAY

#HIGH TEMPS
sat_value_alb = str(data.iloc[0, 8])
sat_value_buf = str(data.iloc[7, 8])
sat_value_syr = str(data.iloc[39, 8])
sat_value_bos = str(data.iloc[4, 8])
sat_value_nyc = str(data.iloc[27, 8])
sat_value_phl = str(data.iloc[32, 8])
sat_value_col = str(data.iloc[13, 8])
sat_value_was = str(data.iloc[41, 8])

sun_value_alb = str(data.iloc[0, 10])
sun_value_buf = str(data.iloc[7, 10])
sun_value_syr = str(data.iloc[39, 10])
sun_value_bos = str(data.iloc[4, 10])
sun_value_nyc = str(data.iloc[27, 10])
sun_value_phl = str(data.iloc[32, 10])
sun_value_col = str(data.iloc[13, 10])
sun_value_was = str(data.iloc[41, 10])

sat_value_sac = str(data.iloc[35, 8]) 
sat_value_lax = str(data.iloc[22, 8])
sat_value_las = str(data.iloc[20, 8])
sat_value_phx = str(data.iloc[33, 8])
sat_value_san = str(data.iloc[36, 8])

sun_value_sac = str(data.iloc[35, 10]) 
sun_value_lax = str(data.iloc[22, 10])
sun_value_las = str(data.iloc[20, 10])
sun_value_phx = str(data.iloc[33, 10])
sun_value_san = str(data.iloc[36, 10])

sat_value_min = str(data.iloc[25, 8]) 
sat_value_des = str(data.iloc[15, 8])
sat_value_chi = str(data.iloc[11, 8])
sat_value_grr = str(data.iloc[18, 8])
sat_value_det = str(data.iloc[16, 8])
sat_value_cinn = str(data.iloc[12, 8])
sat_value_clt = str(data.iloc[10, 8])
sat_value_atl = str(data.iloc[1, 8])
sat_value_wil = str(data.iloc[42, 8])
sat_value_char = str(data.iloc[9, 8])
sat_value_sav = str(data.iloc[38, 8])
sat_value_mem = str(data.iloc[23, 8])
sat_value_dal = str(data.iloc[14, 8])
sat_value_new = str(data.iloc[26, 8])
sat_value_tam = str(data.iloc[40, 8])
sat_value_mia = str(data.iloc[24, 8])

sun_value_min = str(data.iloc[25, 10]) 
sun_value_des = str(data.iloc[15, 10])
sun_value_chi = str(data.iloc[11, 10])
sun_value_grr = str(data.iloc[18, 10])
sun_value_det = str(data.iloc[16, 10])
sun_value_cinn = str(data.iloc[12, 10])
sun_value_clt = str(data.iloc[10, 10])
sun_value_atl = str(data.iloc[1, 10])
sun_value_wil = str(data.iloc[42, 10])
sun_value_char = str(data.iloc[9, 10])
sun_value_sav = str(data.iloc[38, 10])
sun_value_mem = str(data.iloc[23, 10])
sun_value_dal = str(data.iloc[14, 10])
sun_value_new = str(data.iloc[26, 10])
sun_value_tam = str(data.iloc[40, 10])
sun_value_mia = str(data.iloc[24, 10])


#WEATHER CONDITIONS
alb_weather_sat = str(data.iloc[0,9])
buf_weather_sat  = str(data.iloc[7,9])
syr_weather_sat  = str(data.iloc[39,9])
bos_weather_sat  = str(data.iloc[4,9])
nyc_weather_sat  = str(data.iloc[27,9])
phl_weather_sat  = str(data.iloc[32,9])
col_weather_sat  = str(data.iloc[13,9])
was_weather_sat  = str(data.iloc[41,9])

alb_weather_sun = str(data.iloc[0,11])
buf_weather_sun = str(data.iloc[7,11])
syr_weather_sun = str(data.iloc[39,11])
bos_weather_sun = str(data.iloc[4,11])
nyc_weather_sun = str(data.iloc[27,11])
phl_weather_sun = str(data.iloc[32,11])
col_weather_sun = str(data.iloc[13,11])
was_weather_sun = str(data.iloc[41,11])

los_weather_sat = str(data.iloc[22,9])
san_weather_sat = str(data.iloc[36,9])
sac_weather_sat = str(data.iloc[35,9])
las_weather_sat = str(data.iloc[20,9])
phx_weather_sat = str(data.iloc[33,9])

los_weather_sun = str(data.iloc[22,11])
san_weather_sun = str(data.iloc[36,11])
sac_weather_sun = str(data.iloc[35,11])
las_weather_sun = str(data.iloc[20,11])
phx_weather_sun = str(data.iloc[33,11])

cin_weather_sat = str(data.iloc[12,9])
des_weather_sat = str(data.iloc[15,9])
min_weather_sat = str(data.iloc[25,9])
chi_weather_sat = str(data.iloc[11,9])
grr_weather_sat = str(data.iloc[18,9])
det_weather_sat = str(data.iloc[16,9])
sav_weather_sat = str(data.iloc[38,9])
char_weather_sat = str(data.iloc[9,9])
wil_weather_sat = str(data.iloc[42,9])
clt_weather_sat = str(data.iloc[10,9])
atl_weather_sat = str(data.iloc[1,9])
tam_weather_sat = str(data.iloc[40,9])
mem_weather_sat = str(data.iloc[23,9])
dal_weather_sat = str(data.iloc[14,9])
new_weather_sat = str(data.iloc[26,9])
mia_weather_sat = str(data.iloc[24,9])

cin_weather_sun = str(data.iloc[12,11])
des_weather_sun = str(data.iloc[15,11])
min_weather_sun = str(data.iloc[25,11])
chi_weather_sun = str(data.iloc[11,11])
grr_weather_sun = str(data.iloc[18,11])
det_weather_sun = str(data.iloc[16,11])
sav_weather_sun = str(data.iloc[38,11])
char_weather_sun = str(data.iloc[9,11])
wil_weather_sun = str(data.iloc[42,11])
clt_weather_sun = str(data.iloc[10,11])
atl_weather_sun = str(data.iloc[1,11])
tam_weather_sun = str(data.iloc[40,11])
mem_weather_sun = str(data.iloc[23,11])
dal_weather_sun = str(data.iloc[14,11])
new_weather_sun = str(data.iloc[26,11])
mia_weather_sun = str(data.iloc[24,11])



# Define the base directory for the image files
base_directory = "/Users/erick/Desktop/Weather_Graphics/Simple Weather Icons/weather_icons/"
base_directory2 = "/Users/erick/Desktop/Weather_Graphics/Simple Weather Icons/weather_icons_night/"

# Define the dictionary mapping weather values to image file paths
weather_image_mapping2 = {
    "sky is clear": "Moon + Stars.png",
    "moderate rain": "Rain.png",
    "light rain": "Rain.png",
    "overcast clouds": "Cloud.png",
    "scattered clouds": "Night + Clouds.png",
    "broken clouds": "Night + Clouds.png",
    "few clouds": "Moon + Stars.png",
    "heavy intensity rain": "Thunderstorm 2.png",
    "clear sky": "Moon + Stars.png",
    "partly cloudy": "Night + Clouds.png",
    "sunny": "Moon + Stars.png",
    "patchy rain possible": "Rain.png",
    "heavy rain": "Thunderstorm 2.png",
    "thunderstorm with rain": "Thunderstorm 2.png",
    "thunderstorm with heavy rain": "Thunderstorm 2.png",
    "drizzle": "Rain.png",
    "light shower rain": "Rain.png"

}

# Define the dictionary mapping weather values to image file paths
weather_image_mapping = {
    "sky is clear": "Sun 3.png",
    "moderate rain": "Rain.png",
    "light rain": "Rain + Sun.png",
    "overcast clouds": "Cloud.png",
    "scattered clouds": "Sun & Clouds.png",
    "broken clouds": "Sun & Clouds.png",
    "few clouds": "Sun 3.png",
    "heavy intensity rain": "Thunderstorm & Sun.png",
    "clear sky": "Sun 3.png",
    "partly cloudy": "Sun & Clouds.png",
    "sunny": "Sun 3.png",
    "patchy rain possible": "Rain + Sun.png",
    "heavy rain": "Thunderstorm & Sun.png",
    "thunderstorm with rain": "Thunderstorm & Sun.png",
    "thunderstorm with heavy rain": "Thunderstorm 2.png",
    "drizzle": "Rain.png",
    "light shower rain": "Rain.png"
}


presentation = Presentation(powerpoint_file_path)


#WEATHER GRAPHIC CREATION BEGINS HERE *****************************************
#SATURDAY ***********************************************************************

slide = presentation.slides[slide_index]

albany = 29  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox = slide.shapes[albany].text_frame

buffalo = 12  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox2 = slide.shapes[buffalo].text_frame

syracuse = 28  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox3 = slide.shapes[syracuse].text_frame

boston = 30  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox4 = slide.shapes[boston].text_frame

newyork = 8  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox5 = slide.shapes[newyork].text_frame

philadelphia = 24  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox6 = slide.shapes[philadelphia].text_frame

columbus = 16  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox7 = slide.shapes[columbus].text_frame

washington = 20  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox8 = slide.shapes[washington].text_frame

#CLEAR TEXT, ADD NEW VALUES
textbox.clear()
textbox.text = sat_value_alb
textbox2.clear()
textbox2.text = sat_value_buf
textbox3.clear()
textbox3.text = sat_value_syr
textbox4.clear()
textbox4.text = sat_value_bos
textbox5.clear()
textbox5.text = sat_value_nyc
textbox6.clear()
textbox6.text = sat_value_phl
textbox7.clear()
textbox7.text = sat_value_col
textbox8.clear()
textbox8.text = sat_value_was

#FORMATTING NEW TEXT

for paragraph in textbox.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox2.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox3.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox4.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox5.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER


for paragraph in textbox6.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox7.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox8.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

# Assign image file paths based on weather values using the dictionary mapping
nyc_image_file = base_directory + weather_image_mapping.get(nyc_weather_sat, "Wind.png")
alb_image_file = base_directory + weather_image_mapping.get(alb_weather_sat, "Wind.png")
bos_image_file = base_directory + weather_image_mapping.get(bos_weather_sat, "Wind.png")
buf_image_file = base_directory + weather_image_mapping.get(buf_weather_sat, "Wind.png")
syr_image_file = base_directory + weather_image_mapping.get(syr_weather_sat, "Wind.png")
phl_image_file = base_directory + weather_image_mapping.get(phl_weather_sat, "Wind.png")
was_image_file = base_directory + weather_image_mapping.get(was_weather_sat, "Wind.png")
col_image_file = base_directory + weather_image_mapping.get(col_weather_sat, "Wind.png")


# Iterate through the slides and replace images
for slide in presentation.slides:
    for shape in slide.shapes:
        if shape.name == "nyc_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(nyc_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "alb_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(alb_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "bos_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(bos_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "buf_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(buf_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "syr_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(syr_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "phl_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(phl_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "was_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(was_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "col_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(col_image_file, shape.left, shape.top, shape.width, shape.height)


#SUNDAY ***********************************************************************

slide = presentation.slides[slide_index2]

albany2 = 29  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textboxa = slide.shapes[albany2].text_frame

buffalo2 = 12  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox2a = slide.shapes[buffalo2].text_frame

syracuse2 = 28  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox3a = slide.shapes[syracuse2].text_frame

boston2 = 30  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox4a = slide.shapes[boston2].text_frame

newyork2 = 8  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox5a = slide.shapes[newyork2].text_frame

philadelphia2 = 24  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox6a = slide.shapes[philadelphia2].text_frame

columbus2 = 16  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox7a = slide.shapes[columbus2].text_frame

washington2 = 20  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox8a = slide.shapes[washington2].text_frame

#CLEAR TEXT, ADD NEW VALUES
textboxa.clear()
textboxa.text = sun_value_alb
textbox2a.clear()
textbox2a.text = sun_value_buf
textbox3a.clear()
textbox3a.text = sun_value_syr
textbox4a.clear()
textbox4a.text = sun_value_bos
textbox5a.clear()
textbox5a.text = sun_value_nyc
textbox6a.clear()
textbox6a.text = sun_value_phl
textbox7a.clear()
textbox7a.text = sun_value_col
textbox8a.clear()
textbox8a.text = sun_value_was

#FORMATTING NEW TEXT

for paragraph in textboxa.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox2a.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox3a.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox4a.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox5a.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER


for paragraph in textbox6a.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox7a.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox8a.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

# Assign image file paths based on weather values using the dictionary mapping
nyc_image_file2 = base_directory + weather_image_mapping.get(nyc_weather_sun, "Wind.png")
alb_image_file2 = base_directory + weather_image_mapping.get(alb_weather_sun, "Wind.png")
bos_image_file2 = base_directory + weather_image_mapping.get(bos_weather_sun, "Wind.png")
buf_image_file2 = base_directory + weather_image_mapping.get(buf_weather_sun, "Wind.png")
syr_image_file2 = base_directory + weather_image_mapping.get(syr_weather_sun, "Wind.png")
phl_image_file2 = base_directory + weather_image_mapping.get(phl_weather_sun, "Wind.png")
was_image_file2 = base_directory + weather_image_mapping.get(was_weather_sun, "Wind.png")
col_image_file2 = base_directory + weather_image_mapping.get(col_weather_sun, "Wind.png")


# Iterate through the slides and replace images
for slide in presentation.slides:
    for shape in slide.shapes:
        if shape.name == "nyc_icon2":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(nyc_image_file2, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "alb_icon2":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(alb_image_file2, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "bos_icon2":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(bos_image_file2, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "buf_icon2":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(buf_image_file2, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "syr_icon2":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(syr_image_file2, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "phl_icon2":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(phl_image_file2, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "was_icon2":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(was_image_file2, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "col_icon2":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(col_image_file2, shape.left, shape.top, shape.width, shape.height)


#(((((((((((((((((((((((((((((SLIDES FOR WEST COAST WEEKEND)))))))))))))))))))))))))))))))))))#

##### SATURDAY ******************************************************
slide = presentation.slides[slide_index3]

sacramento = 5  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox = slide.shapes[sacramento].text_frame

losangeles = 9  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox2 = slide.shapes[losangeles].text_frame

lasvegas = 17  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox3 = slide.shapes[lasvegas].text_frame


sandiego = 13  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox4 = slide.shapes[sandiego].text_frame

phoenix = 21  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox5 = slide.shapes[phoenix].text_frame


#CLEAR TEXT, ADD NEW VALUES
textbox.clear()
textbox.text = sat_value_sac
textbox2.clear()
textbox2.text = sat_value_lax
textbox3.clear()
textbox3.text = sat_value_las
textbox4.clear()
textbox4.text = sat_value_san
textbox5.clear()
textbox5.text = sat_value_phx

#FORMATTING NEW TEXT


for paragraph in textbox.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox2.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox3.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER


for paragraph in textbox4.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER


for paragraph in textbox5.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

# Assign image file paths based on weather values using the dictionary mapping
los_image_file = base_directory + weather_image_mapping.get(los_weather_sat, "Wind.png")
san_image_file = base_directory + weather_image_mapping.get(san_weather_sat, "Wind.png")
sac_image_file = base_directory + weather_image_mapping.get(sac_weather_sat, "Wind.png")
las_image_file = base_directory + weather_image_mapping.get(las_weather_sat, "Wind.png")
phx_image_file = base_directory + weather_image_mapping.get(phx_weather_sat, "Wind.png")

for slide in presentation.slides:
    for shape in slide.shapes:
        if shape.name == "los_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(los_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "san_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(san_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "sac_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(sac_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "las_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(las_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "phx_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(phx_image_file, shape.left, shape.top, shape.width, shape.height)

##### SUNDAY WEST COAST ******************************************************
slide = presentation.slides[slide_index4]

sacramento = 5  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox = slide.shapes[sacramento].text_frame

losangeles = 9  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox2 = slide.shapes[losangeles].text_frame

lasvegas = 17  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox3 = slide.shapes[lasvegas].text_frame


sandiego = 13  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox4 = slide.shapes[sandiego].text_frame

phoenix = 21  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox5 = slide.shapes[phoenix].text_frame


#CLEAR TEXT, ADD NEW VALUES
textbox.clear()
textbox.text = sun_value_sac
textbox2.clear()
textbox2.text = sun_value_lax
textbox3.clear()
textbox3.text = sun_value_las
textbox4.clear()
textbox4.text = sun_value_san
textbox5.clear()
textbox5.text = sun_value_phx

#FORMATTING NEW TEXT


for paragraph in textbox.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox2.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox3.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER


for paragraph in textbox4.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER


for paragraph in textbox5.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

# Assign image file paths based on weather values using the dictionary mapping
los_image_file = base_directory + weather_image_mapping.get(los_weather_sun, "Wind.png")
san_image_file = base_directory + weather_image_mapping.get(san_weather_sun, "Wind.png")
sac_image_file = base_directory + weather_image_mapping.get(sac_weather_sun, "Wind.png")
las_image_file = base_directory + weather_image_mapping.get(las_weather_sun, "Wind.png")
phx_image_file = base_directory + weather_image_mapping.get(phx_weather_sun, "Wind.png")

for slide in presentation.slides:
    for shape in slide.shapes:
        if shape.name == "los_icon2":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(los_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "san_ico2n":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(san_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "sac_icon2":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(sac_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "las_icon2":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(las_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "phx_icon2":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(phx_image_file, shape.left, shape.top, shape.width, shape.height)


#(((((((((((((((((((((((((((((SLIDES FOR MIDWEST WEEKEND)))))))))))))))))))))))))))))))))))#

#******** SATURDAY ********#

slide = presentation.slides[slide_index5]

cin_image_file = base_directory + weather_image_mapping.get(cin_weather_sat, "Wind.png")
des_image_file = base_directory + weather_image_mapping.get(des_weather_sat, "Wind.png")
min_image_file = base_directory + weather_image_mapping.get(min_weather_sat, "Wind.png")
chi_image_file = base_directory + weather_image_mapping.get(chi_weather_sat, "Wind.png")
grr_image_file = base_directory + weather_image_mapping.get(grr_weather_sat, "Wind.png")
det_image_file = base_directory + weather_image_mapping.get(det_weather_sat, "Wind.png")
sav_image_file = base_directory + weather_image_mapping.get(sav_weather_sat, "Wind.png")
char_image_file = base_directory + weather_image_mapping.get(char_weather_sat, "Wind.png")
wil_image_file = base_directory + weather_image_mapping.get(wil_weather_sat, "Wind.png")
clt_image_file = base_directory + weather_image_mapping.get(clt_weather_sat, "Wind.png")
atl_image_file = base_directory + weather_image_mapping.get(atl_weather_sat, "Wind.png")
tam_image_file = base_directory + weather_image_mapping.get(tam_weather_sat, "Wind.png")
mem_image_file = base_directory + weather_image_mapping.get(mem_weather_sat, "Wind.png")
dal_image_file = base_directory + weather_image_mapping.get(dal_weather_sat, "Wind.png")
new_image_file = base_directory + weather_image_mapping.get(new_weather_sat, "Wind.png")
mia_image_file = base_directory + weather_image_mapping.get(mia_weather_sat, "Wind.png")

cin_image_file2 = base_directory + weather_image_mapping.get(cin_weather_sun, "Wind.png")
des_image_file2 = base_directory + weather_image_mapping.get(des_weather_sun, "Wind.png")
min_image_file2 = base_directory + weather_image_mapping.get(min_weather_sun, "Wind.png")
chi_image_file2 = base_directory + weather_image_mapping.get(chi_weather_sun, "Wind.png")
grr_image_file2 = base_directory + weather_image_mapping.get(grr_weather_sun, "Wind.png")
det_image_file2 = base_directory + weather_image_mapping.get(det_weather_sun, "Wind.png")
sav_image_file2 = base_directory + weather_image_mapping.get(sav_weather_sun, "Wind.png")
char_image_file2 = base_directory + weather_image_mapping.get(char_weather_sun, "Wind.png")
wil_image_file2 = base_directory + weather_image_mapping.get(wil_weather_sun, "Wind.png")
clt_image_file2 = base_directory + weather_image_mapping.get(clt_weather_sun, "Wind.png")
atl_image_file2 = base_directory + weather_image_mapping.get(atl_weather_sun, "Wind.png")
tam_image_file2 = base_directory + weather_image_mapping.get(tam_weather_sun, "Wind.png")
mem_image_file2 = base_directory + weather_image_mapping.get(mem_weather_sun, "Wind.png")
dal_image_file2 = base_directory + weather_image_mapping.get(dal_weather_sun, "Wind.png")
new_image_file2 = base_directory + weather_image_mapping.get(new_weather_sun, "Wind.png")
mia_image_file2 = base_directory + weather_image_mapping.get(mia_weather_sun, "Wind.png")


minneanpolis = 17  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textboxa = slide.shapes[minneanpolis].text_frame

desmoines = 21  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textboxb = slide.shapes[desmoines].text_frame

grandrapids = 9  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textboxc = slide.shapes[grandrapids].text_frame

chicago = 13  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textboxd = slide.shapes[chicago].text_frame

detroit = 5  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textboxe = slide.shapes[detroit].text_frame

cincinnati = 25  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textboxf = slide.shapes[cincinnati].text_frame

#CLEAR TEXT, ADD NEW VALUES
textboxa.clear()
textboxa.text = sat_value_min
textboxb.clear()
textboxb.text = sat_value_des
textboxc.clear()
textboxc.text = sat_value_grr
textboxd.clear()
textboxd.text = sat_value_chi
textboxe.clear()
textboxe.text = sat_value_det
textboxf.clear()
textboxf.text = sat_value_cinn

for paragraph in textboxa.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textboxb.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textboxc.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textboxd.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textboxe.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textboxf.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

#**********(((((((((( SUNDAY )))))))))) *****************

slide = presentation.slides[slide_index6]

minneanpolis = 16  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textboxa2 = slide.shapes[minneanpolis].text_frame

desmoines = 20  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textboxb2 = slide.shapes[desmoines].text_frame

grandrapids = 8  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textboxc2 = slide.shapes[grandrapids].text_frame

chicago = 12  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textboxd2 = slide.shapes[chicago].text_frame

detroit = 4  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textboxe2 = slide.shapes[detroit].text_frame

cincinnati = 24  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textboxf2 = slide.shapes[cincinnati].text_frame

#CLEAR TEXT, ADD NEW VALUES
textboxa2.clear()
textboxa2.text = sun_value_min
textboxb2.clear()
textboxb2.text = sun_value_des
textboxc2.clear()
textboxc2.text = sun_value_grr
textboxd2.clear()
textboxd2.text = sun_value_chi
textboxe2.clear()
textboxe2.text = sun_value_det
textboxf2.clear()
textboxf2.text = sun_value_cinn

for paragraph in textboxa2.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textboxb2.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textboxc2.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textboxd2.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textboxe2.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textboxf2.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

######################## DEEP SOUTH WEEKEND ################################

#####***((((((((((((((((((SATURDAY))))))))))))))))))***#########

slide = presentation.slides[slide_index7]

clt = 28  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox14 = slide.shapes[clt].text_frame

atl = 8  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox15 = slide.shapes[atl].text_frame

sav = 32  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox16 = slide.shapes[sav].text_frame

mem = 20  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox17 = slide.shapes[mem].text_frame

dal = 16  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox18 = slide.shapes[dal].text_frame

new = 12  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox19 = slide.shapes[new].text_frame

tam = 24  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox20 = slide.shapes[tam].text_frame

mia = 4  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox21 = slide.shapes[mia].text_frame

#CLEAR TEXT, ADD NEW VALUES
textbox14.clear()
textbox14.text = sat_value_clt
textbox15.clear()
textbox15.text = sat_value_atl
textbox16.clear()
textbox16.text = sat_value_sav
textbox17.clear()
textbox17.text = sat_value_mem
textbox18.clear()
textbox18.text = sat_value_dal
textbox19.clear()
textbox19.text = sat_value_new
textbox20.clear()
textbox20.text = sat_value_tam
textbox21.clear()
textbox21.text = sat_value_mia

#FORMATTING NEW TEXT
for paragraph in textbox14.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox15.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox16.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox17.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox18.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER


for paragraph in textbox19.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox20.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox21.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

#####***((((((((((((((((((SUNDAY))))))))))))))))))***#########

slide = presentation.slides[slide_index7]

clt = 28  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox14a = slide.shapes[clt].text_frame

atl = 8  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox15a = slide.shapes[atl].text_frame

sav = 32  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox16a = slide.shapes[sav].text_frame

mem = 20  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox17a = slide.shapes[mem].text_frame

dal = 16  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox18a = slide.shapes[dal].text_frame

new = 12  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox19a = slide.shapes[new].text_frame

tam = 24  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox20a = slide.shapes[tam].text_frame

mia = 4  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox21a = slide.shapes[mia].text_frame

#CLEAR TEXT, ADD NEW VALUES
textbox14a.clear()
textbox14a.text = sun_value_clt
textbox15a.clear()
textbox15a.text = sun_value_atl
textbox16a.clear()
textbox16a.text = sun_value_sav
textbox17a.clear()
textbox17a.text = sun_value_mem
textbox18a.clear()
textbox18a.text = sun_value_dal
textbox19a.clear()
textbox19a.text = sun_value_new
textbox20a.clear()
textbox20a.text = sun_value_tam
textbox21a.clear()
textbox21a.text = sun_value_mia

#FORMATTING NEW TEXT
for paragraph in textbox14a.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox15a.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox16a.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox17a.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox18a.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER


for paragraph in textbox19a.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox20a.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox21a.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER


for slide in presentation.slides:
    for shape in slide.shapes:
        if shape.name == "cin_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(cin_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "des_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(des_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "min_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(min_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "chi_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(chi_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "grr_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(grr_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "det_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(det_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "sav_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(sav_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "char_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(char_image_file, shape.left, shape.top, shape.width, shape.height)        
        elif shape.name == "wil_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(wil_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "clt_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(clt_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "tam_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(tam_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "mem_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(mem_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "dal_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(dal_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "new_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(new_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "mia_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(mia_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "atl_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(atl_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "cin_icon2":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(cin_image_file2, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "des_icon2":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(des_image_file2, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "min_icon2":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(min_image_file2, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "chi_icon2":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(chi_image_file2, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "grr_icon2":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(grr_image_file2, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "det_icon2":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(det_image_file2, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "sav_icon2":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(sav_image_file2, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "char_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(char_image_file2, shape.left, shape.top, shape.width, shape.height)        
        elif shape.name == "wil_icon2":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(wil_image_file2, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "clt_icon2":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(clt_image_file2, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "tam_icon2":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(tam_image_file2, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "mem_icon2":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(mem_image_file2, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "dal_icon2":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(dal_image_file2, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "new_icon2":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(new_image_file2, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "mia_icon2":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(mia_image_file2, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "atl_icon2":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(atl_image_file2, shape.left, shape.top, shape.width, shape.height)

#THIS UPDATES THE PRESENTATION
updated_powerpoint_file_path = "/Users/erick/Desktop/weekend_weather.pptx"
presentation.save(updated_powerpoint_file_path)