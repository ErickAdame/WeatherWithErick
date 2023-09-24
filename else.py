
import pandas as pd
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

from pptx.compat import BytesIO
from pptx.enum.shapes import PP_PLACEHOLDER, PROG_ID
from pptx.media import SPEAKER_IMAGE_BYTES, Video
from pptx.opc.constants import CONTENT_TYPE as CT
from pptx.oxml.ns import qn
from pptx.oxml.shapes.graphfrm import CT_GraphicalObjectFrame
from pptx.oxml.shapes.picture import CT_Picture
from pptx.oxml.simpletypes import ST_Direction
from pptx.shapes.autoshape import AutoShapeType, Shape
from pptx.shapes.base import BaseShape
from pptx.shapes.connector import Connector
from pptx.shapes.freeform import FreeformBuilder
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.group import GroupShape
from pptx.shapes.picture import Movie, Picture
from pptx.enum.text import PP_ALIGN





# Set the file paths
powerpoint_file_path = "/Users/erick/Desktop/Graphics_Templates/template_elsewhere.pptx"
csv_file_path = "/Users/erick/Desktop/city_high_and Lows.csv"
csv_file_pathb = "/Users/erick/Desktop/city_high_and Lows_added.csv"
csv_file_path2 = "/Users/erick/Desktop/grand rapids_7_day_forecast.csv"
csv_file_path3 = "/Users/erick/Desktop/day_part_data.csv"
csv_file_path4 = "/Users/erick/Desktop/fort lauderdale_7_day_forecast.csv"
csv_file_path5 = "/Users/erick/Desktop/charlotte_7_day_forecast.csv"



# Set the slide index and text box index of the PowerPoint slide to update
 # REMINDER: Slide index is 0-based, so slide 7 corresponds to index 6
slide_index = 4 #Great Lakes Forecast
slide_index2 = 9 #Southeast Forecast
slide_index3 = 10 #Deep South Forecast
slide_index4 = 6 #Grand Rapids Daypart
slide_index5 = 12 #Fort Lauderdale Day Part
slide_index6 = 7 #Grand Rapids 7 Day
slide_index7 = 8 #Chicago 7 Day
slide_index8 = 11 #Charlotte 7 Day
slide_index9 = 13 #Fort Lauderdale 7 Day
slide_index10 = 5 #Ohio River Forecast

#DATA SOURCE
data = pd.read_csv(csv_file_path)
data2 = pd.read_csv(csv_file_path2)
data3 = pd.read_csv(csv_file_path3)
data4 = pd.read_csv(csv_file_path4)
data5 = pd.read_csv(csv_file_path5)
data6 = pd.read_csv(csv_file_pathb)


#ASSIGN ALL CELL VALUES

#HIGH TEMPS
cell_value_min = str(data.iloc[25, 2]) 
cell_value_des = str(data.iloc[15, 2])
cell_value_chi = str(data.iloc[11, 2])
cell_value_grr = str(data2.iloc[0, 2])
cell_value_det = str(data.iloc[16, 2])
cell_value_cinn = str(data.iloc[12, 2])
cell_value_clt = str(data5.iloc[0, 2])
cell_value_atl = str(data.iloc[1, 2])
cell_value_wil = str(data.iloc[42, 2])
cell_value_char = str(data.iloc[9, 2])
cell_value_sav = str(data.iloc[38, 2])
cell_value_mem = str(data.iloc[23, 2])
cell_value_dal = str(data.iloc[14, 2])
cell_value_new = str(data.iloc[26, 2])
cell_value_tam = str(data.iloc[40, 2])
cell_value_mia = str(data.iloc[24, 2])
cell_value_pitt = str(data6.iloc[3, 2])
cell_value_roa = str(data6.iloc[4, 2])
cell_value_nash = str(data6.iloc[2, 2])
cell_value_lou = str(data6.iloc[1, 2])
cell_value_stl = str(data6.iloc[5, 2])

#DAYPART TEMPS
daypart1_value = str(data3.iloc[5, 1])
daypart2_value = str(data2.iloc[0, 2])
daypart3_value = str(data3.iloc[5, 6])

daypart4_value = str(data3.iloc[7, 1])
daypart5_value = str(data4.iloc[0, 2])
daypart6_value = str(data3.iloc[7, 6])

#WEATHER CONDITIONS
cin_weather = str(data.iloc[12,4])
des_weather = str(data.iloc[15,4])
min_weather = str(data.iloc[25,4])
chi_weather = str(data.iloc[11,4])
grr_weather = str(data.iloc[18,4])
det_weather = str(data.iloc[16,4])
sav_weather = str(data.iloc[38,4])
char_weather = str(data.iloc[9,4])
wil_weather = str(data.iloc[42,4])
clt_weather = str(data.iloc[10,4])
atl_weather = str(data.iloc[1,4])
tam_weather = str(data.iloc[40,4])
mem_weather = str(data.iloc[23,4])
dal_weather = str(data.iloc[14,4])
new_weather = str(data.iloc[26,4])
mia_weather = str(data.iloc[24,4])
pitt_weather = str(data6.iloc[3,4])
roa_weather = str(data6.iloc[4,4])
lou_weather = str(data6.iloc[1,4])
nash_weather = str(data6.iloc[2,4])
stl_weather = str(data6.iloc[5,4])

#DAYPART WEATHER
daypart1_weather = str(data3.iloc[4,1])
daypart2_weather = str(data3.iloc[4,4])
daypart3_weather = str(data3.iloc[4,6])

daypart4_weather = str(data3.iloc[6,1])
daypart5_weather = str(data3.iloc[6,4])
daypart6_weather = str(data3.iloc[6,6])



#IMAGE MAPPING
# Define the base directory for the image files
base_directory = "/Users/erick/Desktop/Weather_Graphics/Simple Weather Icons/weather_icons/"
base_directory2 = "/Users/erick/Desktop/Weather_Graphics/Simple Weather Icons/weather_icons_night/"

weather_image_mapping2 = {
    "sky is clear": "Moon + Stars.png",
    "moderate rain": "Rain.png",
    "light rain": "Rain.png",
    "overcast clouds": "Cloud.png",
    "scattered clouds": "Night + Clouds.png",
    "broken clouds": "Night + Clouds.png",
    "few clouds": "Moon + Stars.png",
    "heavy intensity rain": "Thunderstorm 2.png",
    "clear sky": "Moon + Stars.png",
    "partly cloudy": "Night + Clouds.png",
    "sunny": "Moon + Stars.png",
    "patchy rain possible": "Rain.png",
    "heavy rain": "Thunderstorm 2.png",
    "thunderstorm with rain": "Thunderstorm 2.png",
    "thunderstorm with heavy rain": "Thunderstorm 2.png",
    "drizzle": "Rain.png",
    "light shower rain": "Rain.png"

}

# Define the dictionary mapping weather values to image file paths
weather_image_mapping = {
    "sky is clear": "Sun 3.png",
    "moderate rain": "Rain.png",
    "light rain": "Rain + Sun.png",
    "overcast clouds": "Cloud.png",
    "scattered clouds": "Sun & Clouds.png",
    "broken clouds": "Sun & Clouds.png",
    "few clouds": "Sun 3.png",
    "heavy intensity rain": "Thunderstorm & Sun.png",
    "clear sky": "Sun 3.png",
    "partly cloudy": "Sun & Clouds.png",
    "sunny": "Sun 3.png",
    "patchy rain possible": "Rain + Sun.png",
    "heavy rain": "Thunderstorm & Sun.png",
    "thunderstorm with rain": "Thunderstorm & Sun.png",
    "thunderstorm with heavy rain": "Thunderstorm 2.png",
    "drizzle": "Rain.png",
    "light shower rain": "Rain.png"

    # Add more mappings for other weather conditions
}

# Assign image file paths based on weather values using the dictionary mapping
cin_image_file = base_directory + weather_image_mapping.get(cin_weather, "Wind.png")
des_image_file = base_directory + weather_image_mapping.get(des_weather, "Wind.png")
min_image_file = base_directory + weather_image_mapping.get(min_weather, "Wind.png")
chi_image_file = base_directory + weather_image_mapping.get(chi_weather, "Wind.png")
grr_image_file = base_directory + weather_image_mapping.get(grr_weather, "Wind.png")
det_image_file = base_directory + weather_image_mapping.get(det_weather, "Wind.png")
sav_image_file = base_directory + weather_image_mapping.get(sav_weather, "Wind.png")
char_image_file = base_directory + weather_image_mapping.get(char_weather, "Wind.png")
wil_image_file = base_directory + weather_image_mapping.get(wil_weather, "Wind.png")
clt_image_file = base_directory + weather_image_mapping.get(clt_weather, "Wind.png")
atl_image_file = base_directory + weather_image_mapping.get(atl_weather, "Wind.png")
tam_image_file = base_directory + weather_image_mapping.get(tam_weather, "Wind.png")
mem_image_file = base_directory + weather_image_mapping.get(mem_weather, "Wind.png")
dal_image_file = base_directory + weather_image_mapping.get(dal_weather, "Wind.png")
new_image_file = base_directory + weather_image_mapping.get(new_weather, "Wind.png")
mia_image_file = base_directory + weather_image_mapping.get(mia_weather, "Wind.png")



presentation = Presentation(powerpoint_file_path)



# SLIDE NUMBER 7 BEGINS HERE **************************


slide = presentation.slides[slide_index]


minneanpolis = 17  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textboxa = slide.shapes[minneanpolis].text_frame

desmoines = 21  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textboxb = slide.shapes[desmoines].text_frame

grandrapids = 9  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textboxc = slide.shapes[grandrapids].text_frame

chicago = 13  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textboxd = slide.shapes[chicago].text_frame

detroit = 5  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textboxe = slide.shapes[detroit].text_frame

cincinnati = 25  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textboxf = slide.shapes[cincinnati].text_frame



#CLEAR TEXT, ADD NEW VALUES
textboxa.clear()
textboxa.text = cell_value_min
textboxb.clear()
textboxb.text = cell_value_des
textboxc.clear()
textboxc.text = cell_value_grr
textboxd.clear()
textboxd.text = cell_value_chi
textboxe.clear()
textboxe.text = cell_value_det
textboxf.clear()
textboxf.text = cell_value_cinn


#FORMATTING NEW TEXT


for paragraph in textboxa.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textboxb.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textboxc.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textboxd.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textboxe.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER


for paragraph in textboxf.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER


# SLIDE NUMBER 9 BEGINS HERE ***********************************

slide = presentation.slides[slide_index2]


clt = 9  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox9 = slide.shapes[clt].text_frame

char = 17  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox10 = slide.shapes[char].text_frame

wil = 13  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox11 = slide.shapes[wil].text_frame

atl = 5  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox12 = slide.shapes[atl].text_frame

sav = 21  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox13 = slide.shapes[sav].text_frame


#CLEAR TEXT, ADD NEW VALUES
textbox9.clear()
textbox9.text = cell_value_clt
textbox10.clear()
textbox10.text = cell_value_char
textbox11.clear()
textbox11.text = cell_value_wil
textbox12.clear()
textbox12.text = cell_value_atl
textbox13.clear()
textbox13.text = cell_value_sav


#FORMATTING NEW TEXT

for paragraph in textbox9.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox10.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox11.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox12.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox13.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER



# SLIDE NUMBER 9 BEGINS HERE **************************


slide = presentation.slides[slide_index3]


clt = 29  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox14 = slide.shapes[clt].text_frame

atl = 9  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox15 = slide.shapes[atl].text_frame

sav = 33  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox16 = slide.shapes[sav].text_frame

mem = 21  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox17 = slide.shapes[mem].text_frame

dal = 17  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox18 = slide.shapes[dal].text_frame

new = 13  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox19 = slide.shapes[new].text_frame

tam = 25  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox20 = slide.shapes[tam].text_frame

mia = 5  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox21 = slide.shapes[mia].text_frame


#CLEAR TEXT, ADD NEW VALUES
textbox14.clear()
textbox14.text = cell_value_clt
textbox15.clear()
textbox15.text = cell_value_atl
textbox16.clear()
textbox16.text = cell_value_sav
textbox17.clear()
textbox17.text = cell_value_mem
textbox18.clear()
textbox18.text = cell_value_dal
textbox19.clear()
textbox19.text = cell_value_new
textbox20.clear()
textbox20.text = cell_value_tam
textbox21.clear()
textbox21.text = cell_value_mia

#FORMATTING NEW TEXT


for paragraph in textbox14.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox15.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox16.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox17.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox18.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER


for paragraph in textbox19.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox20.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox21.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER


for slide in presentation.slides:
    for shape in slide.shapes:
        if shape.name == "cin_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(cin_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "des_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(des_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "min_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(min_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "chi_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(chi_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "grr_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(grr_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "det_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(det_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "sav_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(sav_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "char_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(char_image_file, shape.left, shape.top, shape.width, shape.height)        
        elif shape.name == "wil_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(wil_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "clt_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(clt_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "tam_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(tam_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "mem_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(mem_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "dal_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(dal_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "new_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(new_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "mia_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(mia_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "atl_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(atl_image_file, shape.left, shape.top, shape.width, shape.height)

#BEGIN GRR Slide
slide = presentation.slides[slide_index4]

#SELECT TEXT BOXES
daypart1_temp = 11
textbox101 = slide.shapes[daypart1_temp].text_frame
daypart2_temp = 12
textbox102 = slide.shapes[daypart2_temp].text_frame
daypart3_temp = 13
textbox103 = slide.shapes[daypart3_temp].text_frame

#CLEAR AND ADD NEW VALUES
textbox101.clear()
textbox101.text = daypart1_value
textbox102.clear()
textbox102.text = daypart2_value
textbox103.clear()
textbox103.text = daypart3_value


for paragraph in textbox101.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(66)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = False  
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox102.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(66)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = False  
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox103.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(66)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = False  
    paragraph.alignment = PP_ALIGN.CENTER

# Define the base directory for the image files
base_directory = "/Users/erick/Desktop/Weather_Graphics/Simple Weather Icons/weather_icons/"



# Assign image file paths based on weather values using the dictionary mapping
daypart1_image_file = base_directory + weather_image_mapping.get(daypart1_weather, "Wind.png")
daypart2_image_file = base_directory + weather_image_mapping.get(daypart2_weather, "Wind.png")
daypart3_image_file = base_directory2 + weather_image_mapping2.get(daypart3_weather, "Wind.png")



# Iterate through the slides and replace images
for slide in presentation.slides:
    for shape in slide.shapes:
        if shape.name == "daypart1_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(daypart1_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "daypart2_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(daypart2_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "daypart3_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(daypart3_image_file, shape.left, shape.top, shape.width, shape.height)




#BEGIN FLL Slide
slide = presentation.slides[slide_index5]

#SELECT TEXT BOXES
daypart4_temp = 12
textbox104 = slide.shapes[daypart4_temp].text_frame
daypart5_temp = 13
textbox105 = slide.shapes[daypart5_temp].text_frame
daypart6_temp = 14
textbox106 = slide.shapes[daypart6_temp].text_frame


#CLEAR AND ADD NEW VALUES
textbox104.clear()
textbox104.text = daypart4_value
textbox105.clear()
textbox105.text = daypart5_value
textbox106.clear()
textbox106.text = daypart6_value

for paragraph in textbox104.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(66)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = False  
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox105.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(66)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = False  
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox106.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(66)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = False  
    paragraph.alignment = PP_ALIGN.CENTER


# Assign image file paths based on weather values using the dictionary mapping
daypart4_image_file = base_directory + weather_image_mapping.get(daypart4_weather, "Wind.png")
daypart5_image_file = base_directory + weather_image_mapping.get(daypart5_weather, "Wind.png")
daypart6_image_file = base_directory2 + weather_image_mapping2.get(daypart6_weather, "Wind.png")



# Iterate through the slides and replace images
for slide in presentation.slides:
    for shape in slide.shapes:
        if shape.name == "daypart1b":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(daypart4_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "daypart2b_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(daypart5_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "daypart3b_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(daypart6_image_file, shape.left, shape.top, shape.width, shape.height)

slide = presentation.slides[slide_index6]

grr_7day = "/Users/erick/Desktop/grand rapids_7_Day_Forecast.png"
chi_7day = "/Users/erick/Desktop/chicago_7_Day_Forecast.png"
char_7day = "/Users/erick/Desktop/charlotte_7_Day_Forecast.png"
fll_7day  = "/Users/erick/Desktop/fort lauderdale_7_Day_Forecast.png"

# Iterate through the slides and replace images
for slide in presentation.slides:
    for shape in slide.shapes:
        if shape.name == "grr_7day":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(grr_7day, shape.left, shape.top, shape.width, shape.height)  
        
slide = presentation.slides[slide_index7]

# Iterate through the slides and replace images
for slide in presentation.slides:
    for shape in slide.shapes:
        if shape.name == "chi_7day":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(chi_7day, shape.left, shape.top, shape.width, shape.height)

# Iterate through the slides and replace images

slide = presentation.slides[slide_index8]
for slide in presentation.slides:
    for shape in slide.shapes:
         if shape.name == "char_7day":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(char_7day, shape.left, shape.top, shape.width, shape.height)


slide = presentation.slides[slide_index9]
for slide in presentation.slides:
    for shape in slide.shapes:
        if shape.name == "fll_7day":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(fll_7day, shape.left, shape.top, shape.width, shape.height)


#OHIO VALLEY SIDE
slide = presentation.slides[slide_index10]


pitt = 4  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox14 = slide.shapes[pitt].text_frame

cin = 7  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox15 = slide.shapes[cin].text_frame

roa = 19  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox16 = slide.shapes[roa].text_frame

lou = 16  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox17 = slide.shapes[lou].text_frame

nash = 10  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox18 = slide.shapes[nash].text_frame

mem2 = 21  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textboxz = slide.shapes[mem2].text_frame

stl = 13  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox20 = slide.shapes[stl].text_frame


#CLEAR TEXT, ADD NEW VALUES
textbox14.clear()
textbox14.text = cell_value_pitt
textbox15.clear()
textbox15.text = cell_value_cinn
textbox16.clear()
textbox16.text = cell_value_roa
textbox17.clear()
textbox17.text = cell_value_lou
textbox18.clear()
textbox18.text = cell_value_nash
textboxz.clear()
textboxz.text = cell_value_mem
textbox20.clear()
textbox20.text = cell_value_stl

#FORMATTING NEW TEXT


for paragraph in textbox14.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox15.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox16.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox17.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox18.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER


for paragraph in textboxz.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox20.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox21.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

# Assign image file paths based on weather values using the dictionary mapping
cin_image_file = base_directory + weather_image_mapping.get(cin_weather, "Wind.png")
pitt_image_file = base_directory + weather_image_mapping.get(pitt_weather, "Wind.png")
roa_image_file = base_directory + weather_image_mapping.get(roa_weather, "Wind.png")
lou_image_file = base_directory + weather_image_mapping.get(cin_weather, "Wind.png")
nash_image_file = base_directory + weather_image_mapping.get(nash_weather, "Wind.png")
mem_image_file = base_directory + weather_image_mapping.get(mem_weather, "Wind.png")
stl_image_file = base_directory + weather_image_mapping.get(stl_weather, "Wind.png")

for slide in presentation.slides:
    for shape in slide.shapes:
        if shape.name == "cin_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(cin_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "pitt_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(pitt_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "roa_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(roa_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "lou_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(lou_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "nash_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(nash_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "mem_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(mem_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "stl_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(stl_image_file, shape.left, shape.top, shape.width, shape.height)
        
#THIS UPDATES THE PRESENTATION
updated_powerpoint_file_path = "/Users/erick/Desktop/Weather_Elsewhere.pptx"
presentation.save(updated_powerpoint_file_path)

