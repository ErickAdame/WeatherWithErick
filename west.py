import pandas as pd
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

from pptx.compat import BytesIO
from pptx.enum.shapes import PP_PLACEHOLDER, PROG_ID
from pptx.media import SPEAKER_IMAGE_BYTES, Video
from pptx.opc.constants import CONTENT_TYPE as CT
from pptx.oxml.ns import qn
from pptx.oxml.shapes.graphfrm import CT_GraphicalObjectFrame
from pptx.oxml.shapes.picture import CT_Picture
from pptx.oxml.simpletypes import ST_Direction
from pptx.shapes.autoshape import AutoShapeType, Shape
from pptx.shapes.base import BaseShape
from pptx.shapes.connector import Connector
from pptx.shapes.freeform import FreeformBuilder
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.group import GroupShape
from pptx.shapes.picture import Movie, Picture
from pptx.enum.text import PP_ALIGN



# Set the file paths
powerpoint_file_path = "/Users/erick/Desktop/Graphics_Templates/template_west.pptx"
csv_file_path = "/Users/erick/Desktop/city_high_and Lows.csv"
csv_file_path2 = "/Users/erick/Desktop/los angeles_7_day_forecast.csv"
csv_file_path3 = "/Users/erick/Desktop/day_part_data.csv"

# Set the slide index and text box index of the PowerPoint slide to update
 # REMINDER: Slide index is 0-based, so slide 7 corresponds to index 6
slide_index = 7 
slide_index2= 8
slide_index3= 9
slide_index4=1
slide_index5=2

#DATA SOURCE
data = pd.read_csv(csv_file_path)
data2 = pd.read_csv(csv_file_path2)
data3 = pd.read_csv(csv_file_path3)



#ASSIGN ALL CELL VALUES

#HIGH TEMPS
cell_value_sac = str(data.iloc[35, 2])  # Assuming the value is in cell C2 (row 1, column 2)
cell_value_lax = str(data.iloc[22, 2])
cell_value_las = str(data.iloc[20, 2])
cell_value_phx = str(data.iloc[33, 2])
cell_value_san = str(data.iloc[36, 2])
cell_value_bak = str(data.iloc[2, 2])
cell_value_sba = str(data.iloc[37, 2])
cell_value_ocn = str(data.iloc[29, 2])
cell_value_plm = str(data.iloc[30, 2])

#WEATHER CONDITIONS
#WEATHER CONDITIONS
los_weather = str(data.iloc[22,4])
san_weather = str(data.iloc[36,4])
sac_weather = str(data.iloc[35,4])
las_weather = str(data.iloc[20,4])
phx_weather = str(data.iloc[33,4])
bak_weather = str(data.iloc[2,4])
pal_weather = str(data.iloc[30,4])
ocn_weather = str(data.iloc[29,4])
sba_weather = str(data.iloc[37,4])

#DAYPART WEATHER CONDITIONS


#IMAGE MAPPING
# Define the base directory for the image files
base_directory = "/Users/erick/Desktop/Weather_Graphics/Simple Weather Icons/weather_icons/"

# Define the dictionary mapping weather values to image file paths
weather_image_mapping = {
    "sky is clear": "Sun 3.png",
    "moderate rain": "Rain.png",
    "light rain": "Rain + Sun.png",
    "overcast clouds": "Cloud.png",
    "scattered clouds": "Sun & Clouds.png",
    "broken clouds": "Sun & Clouds.png",
    "few clouds": "Sun 3.png",
    "heavy intensity rain": "Thunderstorm & Sun.png",
    "clear sky": "Sun 3.png",
    "partly cloudy": "Sun & Clouds.png",
    "sunny": "Sun 3.png",
    "patchy rain possible": "Rain + Sun.png",
    "heavy rain": "Thunderstorm & Sun.png",
    "thunderstorm with rain": "Thunderstorm & Sun.png",
    "thunderstorm with heavy rain": "Thunderstorm 2.png"

    # Add more mappings for other weather conditions
}

# Assign image file paths based on weather values using the dictionary mapping
los_image_file = base_directory + weather_image_mapping.get(los_weather, "Wind.png")
san_image_file = base_directory + weather_image_mapping.get(san_weather, "Wind.png")
sac_image_file = base_directory + weather_image_mapping.get(sac_weather, "Wind.png")
las_image_file = base_directory + weather_image_mapping.get(las_weather, "Wind.png")
phx_image_file = base_directory + weather_image_mapping.get(phx_weather, "Wind.png")
bak_image_file = base_directory + weather_image_mapping.get(bak_weather, "Wind.png")
pal_image_file = base_directory + weather_image_mapping.get(pal_weather, "Wind.png")
ocn_image_file = base_directory + weather_image_mapping.get(ocn_weather, "Wind.png")
sba_image_file = base_directory + weather_image_mapping.get(sba_weather, "Wind.png")




presentation = Presentation(powerpoint_file_path)


# SLIDE NUMBER 7 BEGINS HERE **************************


slide = presentation.slides[slide_index]


sacramento = 5  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox = slide.shapes[sacramento].text_frame

losangeles = 9  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox2 = slide.shapes[losangeles].text_frame

lasvegas = 17  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox3 = slide.shapes[lasvegas].text_frame


sandiego = 13  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox4 = slide.shapes[sandiego].text_frame

phoenix = 21  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox5 = slide.shapes[phoenix].text_frame



#CLEAR TEXT, ADD NEW VALUES
textbox.clear()
textbox.text = cell_value_sac
textbox2.clear()
textbox2.text = cell_value_lax
textbox3.clear()
textbox3.text = cell_value_las
textbox4.clear()
textbox4.text = cell_value_san
textbox5.clear()
textbox5.text = cell_value_phx

#FORMATTING NEW TEXT


for paragraph in textbox.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox2.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox3.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER


for paragraph in textbox4.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER


for paragraph in textbox5.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER



# SLIDE NUMBER 8 BEGINS HERE ***********************************

slide = presentation.slides[slide_index2]


bakersfield = 5  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox9 = slide.shapes[bakersfield].text_frame

santabarbara = 13  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox10 = slide.shapes[santabarbara].text_frame

oceanside = 9  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox11 = slide.shapes[oceanside].text_frame

palmsprings = 17  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox12 = slide.shapes[palmsprings].text_frame


# In[253]:

#CLEAR TEXT, ADD NEW VALUES
textbox9.clear()
textbox9.text = cell_value_bak
textbox10.clear()
textbox10.text = cell_value_sba
textbox11.clear()
textbox11.text = cell_value_ocn
textbox12.clear()
textbox12.text = cell_value_plm

#FORMATTING NEW TEXT

for paragraph in textbox9.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox10.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox11.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox12.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for slide in presentation.slides:
    for shape in slide.shapes:
        if shape.name == "los_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(los_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "san_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(san_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "sac_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(sac_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "las_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(las_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "phx_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(phx_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "bak_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(bak_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "pal_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(pal_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "sba_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(sba_image_file, shape.left, shape.top, shape.width, shape.height)        
        elif shape.name == "ocn_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(ocn_image_file, shape.left, shape.top, shape.width, shape.height)


# SLIDE NUMBER 16

csv_file_path2 = "/Users/erick/Desktop/los angeles_7_day_forecast.csv"
data2 = pd.read_csv(csv_file_path2)

#ASSIGN HIGH TEMP VALUES
day1_high_value = str(data.iloc[22, 2])
day2_high_value = str(data2.iloc[1, 2])
day3_high_value = str(data2.iloc[2, 2])
day4_high_value = str(data2.iloc[3, 2])
day5_high_value = str(data2.iloc[4, 2])
day6_high_value = str(data2.iloc[5, 2])
day7_high_value = str(data2.iloc[6, 2])

#ASSIGN LOW TEMP VALUES
day1_low_value = str(data2.iloc[1, 3])
day2_low_value = str(data2.iloc[2, 3])
day3_low_value = str(data2.iloc[3, 3])
day4_low_value = str(data2.iloc[4, 3])
day5_low_value = str(data2.iloc[5, 3])
day6_low_value = str(data2.iloc[6, 3])
day7_low_value = str(data2.iloc[6, 3]-1)

slide = presentation.slides[slide_index3]

#SELECT AND MAKE TEXT BOX ACTIVE
day1_high = 9
textbox17 = slide.shapes[day1_high].text_frame
day2_high = 11
textbox18 = slide.shapes[day2_high].text_frame
day3_high = 10
textbox19 = slide.shapes[day3_high].text_frame
day4_high = 12
textbox20 = slide.shapes[day4_high].text_frame
day5_high = 14  
textbox21 = slide.shapes[day5_high].text_frame
day6_high = 13
textbox22 = slide.shapes[day6_high].text_frame
day7_high = 15  
textbox23 = slide.shapes[day7_high].text_frame

day1_low = 16
textbox24 = slide.shapes[day1_low].text_frame
day2_low = 18
textbox25 = slide.shapes[day2_low].text_frame
day3_low = 17
textbox26 = slide.shapes[day3_low].text_frame
day4_low = 19
textbox27 = slide.shapes[day4_low].text_frame
day5_low = 21  
textbox28 = slide.shapes[day5_low].text_frame
day6_low = 20
textbox29 = slide.shapes[day6_low].text_frame
day7_low = 22  
textbox30 = slide.shapes[day7_low].text_frame

#WEATHER CONDITIONS
day1_weather = str(data2.iloc[0, 4]).lower()
day2_weather = str(data2.iloc[1, 4]).lower()
day3_weather = str(data2.iloc[2, 4]).lower()
day4_weather = str(data2.iloc[3, 4]).lower()
day5_weather = str(data2.iloc[4, 4]).lower()
day6_weather = str(data2.iloc[5, 4]).lower()
day7_weather = str(data2.iloc[6, 4]).lower()


# Assign image file paths based on weather values using the dictionary mapping
day1_image_file = base_directory + weather_image_mapping.get(day1_weather, "Wind.png")
day2_image_file = base_directory + weather_image_mapping.get(day2_weather, "Wind.png")
day3_image_file = base_directory + weather_image_mapping.get(day3_weather, "Wind.png")
day4_image_file = base_directory + weather_image_mapping.get(day4_weather, "Wind.png")
day5_image_file = base_directory + weather_image_mapping.get(day5_weather, "Wind.png")
day6_image_file = base_directory + weather_image_mapping.get(day6_weather, "Wind.png")
day7_image_file = base_directory + weather_image_mapping.get(day7_weather, "Wind.png")

#CLEAR TEXT, ADD NEW VALUES
textbox17.clear()
textbox17.text = day1_high_value
textbox18.clear()
textbox18.text = day2_high_value
textbox19.clear()
textbox19.text = day3_high_value
textbox20.clear()
textbox20.text = day4_high_value
textbox21.clear()
textbox21.text = day5_high_value
textbox22.clear()
textbox22.text = day6_high_value
textbox23.clear()
textbox23.text = day7_high_value

textbox24.clear()
textbox24.text = day1_low_value
textbox25.clear()
textbox25.text = day2_low_value
textbox26.clear()
textbox26.text = day3_low_value
textbox27.clear()
textbox27.text = day4_low_value
textbox28.clear()
textbox28.text = day5_low_value
textbox29.clear()
textbox29.text = day6_low_value
textbox30.clear()
textbox30.text = day7_low_value

#FORMATTING NEW TEXT


for paragraph in textbox17.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(50)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = False  
    paragraph.alignment = PP_ALIGN.CENTER


for paragraph in textbox18.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(50)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = False  
    paragraph.alignment = PP_ALIGN.CENTER


for paragraph in textbox19.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(50)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = False 
    paragraph.alignment = PP_ALIGN.CENTER


for paragraph in textbox20.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(50)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = False  
    paragraph.alignment = PP_ALIGN.CENTER


for paragraph in textbox21.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(50)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = False 
    paragraph.alignment = PP_ALIGN.CENTER


for paragraph in textbox22.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(50)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = False  
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox23.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(50)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = False  
    paragraph.alignment = PP_ALIGN.CENTER

# FORMATTING LOW TEMPERATURES


for paragraph in textbox24.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(38)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(180, 199, 231)
        run.font.bold = False  
    paragraph.alignment = PP_ALIGN.CENTER


for paragraph in textbox25.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(38)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(180, 199, 231)
        run.font.bold = False  
    paragraph.alignment = PP_ALIGN.CENTER


for paragraph in textbox26.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(38)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(180, 199, 231)
        run.font.bold = False 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox27.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(38)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(180, 199, 231)
        run.font.bold = False  
    paragraph.alignment = PP_ALIGN.CENTER


for paragraph in textbox28.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(38)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(180, 199, 231)
        run.font.bold = False 
    paragraph.alignment = PP_ALIGN.CENTER


for paragraph in textbox29.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(38)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(180, 199, 231)
        run.font.bold = False  
    paragraph.alignment = PP_ALIGN.CENTER


for paragraph in textbox30.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(38)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(180, 199, 231)
        run.font.bold = False  
    paragraph.alignment = PP_ALIGN.CENTER

# Iterate through the slides and replace images
for slide in presentation.slides:
    for shape in slide.shapes:
        if shape.name == "day1_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(day1_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "day2_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(day2_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "day3_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(day3_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "day4_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(day4_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "day5_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(day5_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "day6_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(day6_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "day7_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(day7_image_file, shape.left, shape.top, shape.width, shape.height)


#SETTINGS FOR FIRST SLIDE

slide = presentation.slides[slide_index4]

csv_file_path3 = "/Users/erick/Desktop/day_part_data.csv"
data3 = pd.read_csv(csv_file_path3)


#ASSIGN HIGH TEMP VALUES
daypart1_value = str(data3.iloc[3, 1])
daypart2_value = str(data.iloc[22, 2])
daypart3_value = str(data3.iloc[3, 5])

#WEATHER CONDITIONS
daypart1_weather = str(data3.iloc[2,2])
daypart2_weather = str(data3.iloc[2,5])
daypart3_weather = str(data3.iloc[2,6])

#SELECT TEXT BOXES
daypart1_temp = 11
textbox101 = slide.shapes[daypart1_temp].text_frame
daypart2_temp = 12
textbox102 = slide.shapes[daypart2_temp].text_frame
daypart3_temp = 13
textbox103 = slide.shapes[daypart3_temp].text_frame

#CLEAR AND ADD NEW VALUES
textbox101.clear()
textbox101.text = daypart1_value
textbox102.clear()
textbox102.text = daypart2_value
textbox103.clear()
textbox103.text = daypart3_value


for paragraph in textbox101.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(66)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = False  
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox102.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(66)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = False  
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox103.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(66)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = False  
    paragraph.alignment = PP_ALIGN.CENTER

# Define the base directory for the image files
base_directory = "/Users/erick/Desktop/Weather_Graphics/Simple Weather Icons/weather_icons/"

base_directory2 = "/Users/erick/Desktop/Weather_Graphics/Simple Weather Icons/weather_icons_night/"

weather_image_mapping2 = {
    "sky is clear": "Moon + Stars.png",
    "moderate rain": "Rain.png",
    "light rain": "Rain.png",
    "overcast clouds": "Cloud.png",
    "scattered clouds": "Night + Clouds.png",
    "broken clouds": "Night + Clouds.png",
    "few clouds": "Moon + Stars.png",
    "heavy intensity rain": "Thunderstorm 2.png",
    "clear sky": "Moon + Stars.png",
    "partly cloudy": "Night + Clouds.png",
    "sunny": "Moon + Stars.png",
    "patchy rain possible": "Rain.png",
    "heavy rain": "Thunderstorm 2.png",
    "thunderstorm with rain": "Thunderstorm 2.png",
    "thunderstorm with heavy rain": "Thunderstorm 2.png"
}

# Assign image file paths based on weather values using the dictionary mapping
daypart1_image_file = base_directory + weather_image_mapping.get(daypart1_weather, "Wind.png")
daypart2_image_file = base_directory + weather_image_mapping.get(daypart2_weather, "Wind.png")
daypart3_image_file = base_directory2 + weather_image_mapping2.get(daypart3_weather, "Wind.png")



# Iterate through the slides and replace images
for slide in presentation.slides:
    for shape in slide.shapes:
        if shape.name == "daypart1_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(daypart1_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "daypart2_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(daypart2_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "daypart3_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(daypart3_image_file, shape.left, shape.top, shape.width, shape.height)

#SETTINGS FOR SECOND SLIDE

slide = presentation.slides[slide_index5]

csv_file_path3 = "/Users/erick/Desktop/day_part_data.csv"
data3 = pd.read_csv(csv_file_path3)


#ASSIGN HIGH TEMP VALUES
daypart4_value = str(data3.iloc[9, 1])
daypart5_value = str(data.iloc[36, 2])
daypart6_value = str(data3.iloc[9, 5])

#WEATHER CONDITIONS
daypart4_weather = str(data3.iloc[7,2])
daypart5_weather = str(data3.iloc[7,5])
daypart6_weather = str(data3.iloc[7,6])

#SELECT TEXT BOXES
daypart4_temp = 11
textbox104 = slide.shapes[daypart1_temp].text_frame
daypart5_temp = 12
textbox105 = slide.shapes[daypart2_temp].text_frame
daypart6_temp = 13
textbox106 = slide.shapes[daypart3_temp].text_frame

#CLEAR AND ADD NEW VALUES
textbox104.clear()
textbox104.text = daypart4_value
textbox105.clear()
textbox105.text = daypart5_value
textbox106.clear()
textbox106.text = daypart6_value


for paragraph in textbox104.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(66)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = False  
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox105.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(66)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = False  
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox106.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(66)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = False  
    paragraph.alignment = PP_ALIGN.CENTER

# Define the base directory for the image files
base_directory = "/Users/erick/Desktop/Weather_Graphics/Simple Weather Icons/weather_icons/"



# Assign image file paths based on weather values using the dictionary mapping
daypart4_image_file = base_directory + weather_image_mapping.get(daypart1_weather, "Wind.png")
daypart5_image_file = base_directory + weather_image_mapping.get(daypart2_weather, "Wind.png")
daypart6_image_file = base_directory2 + weather_image_mapping2.get(daypart3_weather, "Wind.png")



# Iterate through the slides and replace images
for slide in presentation.slides:
    for shape in slide.shapes:
        if shape.name == "daypart1b_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(daypart1_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "daypart2b_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(daypart2_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "daypart3b_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(daypart3_image_file, shape.left, shape.top, shape.width, shape.height)



#THIS UPDATES THE PRESENTATION
updated_powerpoint_file_path = "/Users/erick/Desktop/Weather Gfx WEST.pptx"
presentation.save(updated_powerpoint_file_path)

