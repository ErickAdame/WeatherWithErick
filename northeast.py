

import pandas as pd
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

from pptx.compat import BytesIO
from pptx.enum.shapes import PP_PLACEHOLDER, PROG_ID
from pptx.media import SPEAKER_IMAGE_BYTES, Video
from pptx.opc.constants import CONTENT_TYPE as CT
from pptx.oxml.ns import qn
from pptx.oxml.shapes.graphfrm import CT_GraphicalObjectFrame
from pptx.oxml.shapes.picture import CT_Picture
from pptx.oxml.simpletypes import ST_Direction
from pptx.shapes.autoshape import AutoShapeType, Shape
from pptx.shapes.base import BaseShape
from pptx.shapes.connector import Connector
from pptx.shapes.freeform import FreeformBuilder
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.group import GroupShape
from pptx.shapes.picture import Movie, Picture
from pptx.enum.text import PP_ALIGN



# Set the file paths
powerpoint_file_path = "/Users/erick/Desktop/Graphics_Templates/northeast_template.pptx"
csv_file_path = "/Users/erick/Desktop/city_high_and Lows.csv"

csv_file_path2 = "/Users/erick/Desktop/new york_7_day_forecast.csv"


# Set the slide index and text box index of the PowerPoint slide to update
 # REMINDER: Slide index is 0-based, so slide 7 corresponds to index 6
slide_index = 7 
slide_index2= 8
slide_index3= 17
slide_index4 = 3
slide_index5 = 9
slide_index6 = 18
slide_index7 = 16

#DATA SOURCE
data = pd.read_csv(csv_file_path)
data2 = pd.read_csv(csv_file_path2)



#ASSIGN ALL CELL VALUES

#HIGH TEMPS
cell_value_alb = str(data.iloc[0, 2])  # Assuming the value is in cell C2 (row 1, column 2)
cell_value_buf = str(data.iloc[7, 2])
cell_value_syr = str(data.iloc[39, 2])
cell_value_bos = str(data.iloc[4, 2])
cell_value_nyc = str(data2.iloc[0, 2])
cell_value_phl = str(data.iloc[32, 2])
cell_value_col = str(data.iloc[13, 2])
cell_value_was = str(data.iloc[41, 2])

#LOW TEMPS
cell_value_alb_low = str(data.iloc[0, 3])  # Assuming the value is in cell C2 (row 1, column 2)
cell_value_buf_low = str(data.iloc[7, 3])
cell_value_syr_low = str(data.iloc[39, 3])
cell_value_bos_low = str(data.iloc[4, 3])
cell_value_nyc_low = str(data.iloc[27, 3])
cell_value_phl_low = str(data.iloc[32, 3])
cell_value_col_low = str(data.iloc[13, 3])
cell_value_was_low = str(data.iloc[41, 3])

#APPARENT TEMPS
app_value_alb = str(data.iloc[0, 7])  # Assuming the value is in cell C2 (row 1, column 2)
app_value_buf = str(data.iloc[7, 7])
app_value_syr = str(data.iloc[39, 7])
app_value_bos = str(data.iloc[4, 7])
app_value_nyc = str(data.iloc[27, 7])
app_value_phl = str(data.iloc[32, 7])
app_value_col = str(data.iloc[13, 7])
app_value_was = str(data.iloc[41, 7])

#WEATHER CONDITIONS
alb_weather = str(data.iloc[0,4])
buf_weather = str(data.iloc[7,4])
syr_weather = str(data.iloc[39,4])
bos_weather = str(data.iloc[4,4])
nyc_weather = str(data.iloc[27,4])
phl_weather = str(data.iloc[32,4])
col_weather = str(data.iloc[13,4])
was_weather = str(data.iloc[41,4])

#DAILY RAINFALL
rain_value_alb = str(data.iloc[0, 5])  # Assuming the value is in cell C2 (row 1, column 2)
rain_value_buf = str(data.iloc[7, 5])
rain_value_syr = str(data.iloc[39, 5])
rain_value_bos = str(data.iloc[4, 5])
rain_value_nyc = str(data.iloc[27, 5])
rain_value_phl = str(data.iloc[32, 5])
rain_value_col = str(data.iloc[13, 5])
rain_value_was = str(data.iloc[41, 5])

#WEEKEND DATA
alb_sat = str(data.iloc[0,9])
alb_sun = str(data.iloc[0,11])
bos_sat = str(data.iloc[4,9])
bos_sun = str(data.iloc[4,11])
buf_sat = str(data.iloc[7,9])
buf_sun = str(data.iloc[7,11])
col_sat = str(data.iloc[13,9])
col_sun = str(data.iloc[13,11])
was_sat = str(data.iloc[41,9]) 
was_sun = str(data.iloc[41,11])
phl_sat = str(data.iloc[32,9])
phl_sun = str(data.iloc[32,11])
nyc_sat = str(data.iloc[27,9])
nyc_sun = str(data.iloc[27,11])

alb_sat_temp = str(data.iloc[0,8])
alb_sun_temp = str(data.iloc[0,10])
bos_sat_temp = str(data.iloc[4,8])
bos_sun_temp = str(data.iloc[4,10])
buf_sat_temp = str(data.iloc[7,8])
buf_sun_temp = str(data.iloc[7,10])
col_sat_temp = str(data.iloc[13,8])
col_sun_temp = str(data.iloc[13,10])
was_sat_temp = str(data.iloc[41,8]) 
was_sun_temp = str(data.iloc[41,10])
phl_sat_temp = str(data.iloc[32,8])
phl_sun_temp = str(data.iloc[32,10])
nyc_sat_temp = str(data.iloc[27,8])
nyc_sun_temp = str(data.iloc[27,10])


presentation = Presentation(powerpoint_file_path)


# SLIDE NUMBER 7 BEGINS HERE **************************


slide = presentation.slides[slide_index]


albany = 29  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox = slide.shapes[albany].text_frame

buffalo = 12  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox2 = slide.shapes[buffalo].text_frame

syracuse = 28  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox3 = slide.shapes[syracuse].text_frame

boston = 30  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox4 = slide.shapes[boston].text_frame

newyork = 8  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox5 = slide.shapes[newyork].text_frame

philadelphia = 24  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox6 = slide.shapes[philadelphia].text_frame

columbus = 16  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox7 = slide.shapes[columbus].text_frame

washington = 20  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox8 = slide.shapes[washington].text_frame


# In[253]:

#CLEAR TEXT, ADD NEW VALUES
textbox.clear()
textbox.text = cell_value_alb
textbox2.clear()
textbox2.text = cell_value_buf
textbox3.clear()
textbox3.text = cell_value_syr
textbox4.clear()
textbox4.text = cell_value_bos
textbox5.clear()
textbox5.text = cell_value_nyc
textbox6.clear()
textbox6.text = cell_value_phl
textbox7.clear()
textbox7.text = cell_value_col
textbox8.clear()
textbox8.text = cell_value_was

#FORMATTING NEW TEXT


for paragraph in textbox.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox2.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox3.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox4.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox5.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER


for paragraph in textbox6.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox7.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox8.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

#IMAGE CODING HERE


# Define the base directory for the image files
base_directory = "/Users/erick/Desktop/Weather_Graphics/Simple Weather Icons/weather_icons/"

# Define the dictionary mapping weather values to image file paths
weather_image_mapping = {
    "sky is clear": "Sun 3.png",
    "moderate rain": "Rain.png",
    "light rain": "Rain + Sun.png",
    "overcast clouds": "Cloud.png",
    "scattered clouds": "Sun & Clouds.png",
    "broken clouds": "Sun & Clouds.png",
    "few clouds": "Sun 3.png",
    "heavy intensity rain": "Thunderstorm & Sun.png",
    "clear sky": "Sun 3.png",
    "partly cloudy": "Sun & Clouds.png",
    "sunny": "Sun 3.png",
    "patchy rain possible": "Rain + Sun.png",
    "heavy rain": "Thunderstorm & Sun.png",
    "thunderstorm with rain": "Thunderstorm & Sun.png",
    "thunderstorm with heavy rain": "Thunderstorm 2.png"

    # Add more mappings for other weather conditions
}

# Assign image file paths based on weather values using the dictionary mapping
nyc_image_file = base_directory + weather_image_mapping.get(nyc_weather, "Wind.png")
alb_image_file = base_directory + weather_image_mapping.get(alb_weather, "Wind.png")
bos_image_file = base_directory + weather_image_mapping.get(bos_weather, "Wind.png")
buf_image_file = base_directory + weather_image_mapping.get(buf_weather, "Wind.png")
syr_image_file = base_directory + weather_image_mapping.get(syr_weather, "Wind.png")
phl_image_file = base_directory + weather_image_mapping.get(phl_weather, "Wind.png")
was_image_file = base_directory + weather_image_mapping.get(was_weather, "Wind.png")
col_image_file = base_directory + weather_image_mapping.get(col_weather, "Wind.png")


# Iterate through the slides and replace images
for slide in presentation.slides:
    for shape in slide.shapes:
        if shape.name == "nyc_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(nyc_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "alb_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(alb_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "bos_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(bos_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "buf_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(buf_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "syr_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(syr_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "phl_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(phl_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "was_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(was_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "col_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(col_image_file, shape.left, shape.top, shape.width, shape.height)

# SLIDE NUMBER 8 BEGINS HERE ***********************************

slide = presentation.slides[slide_index2]


albany = 5  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox9 = slide.shapes[albany].text_frame

buffalo = 17  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox10 = slide.shapes[buffalo].text_frame

syracuse = 33  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox11 = slide.shapes[syracuse].text_frame

boston = 9  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox12 = slide.shapes[boston].text_frame

newyork = 13  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox13 = slide.shapes[newyork].text_frame

philadelphia = 29  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox14 = slide.shapes[philadelphia].text_frame

columbus = 21  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox15 = slide.shapes[columbus].text_frame

washington = 25  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox16 = slide.shapes[washington].text_frame


# In[253]:

#CLEAR TEXT, ADD NEW VALUES
textbox9.clear()
textbox9.text = cell_value_alb_low
textbox10.clear()
textbox10.text = cell_value_buf_low
textbox11.clear()
textbox11.text = cell_value_syr_low
textbox12.clear()
textbox12.text = cell_value_bos_low
textbox13.clear()
textbox13.text = cell_value_nyc_low
textbox14.clear()
textbox14.text = cell_value_phl_low
textbox15.clear()
textbox15.text = cell_value_col_low
textbox16.clear()
textbox16.text = cell_value_was_low

#FORMATTING NEW TEXT

for paragraph in textbox9.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox10.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox11.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox12.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox13.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox14.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox15.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox16.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER


# SLIDE NUMBER 16

csv_file_path2 = "/Users/erick/Desktop/new york_7_day_forecast.csv"
data2 = pd.read_csv(csv_file_path2)

#ASSIGN HIGH TEMP VALUES
day1_high_value = str(data2.iloc[0, 2])
day2_high_value = str(data2.iloc[1, 2])
day3_high_value = str(data2.iloc[2, 2])
day4_high_value = str(data2.iloc[3, 2])
day5_high_value = str(data2.iloc[4, 2])
day6_high_value = str(data2.iloc[5, 2])
day7_high_value = str(data2.iloc[6, 2])

#ASSIGN LOW TEMP VALUES
day1_low_value = str(data2.iloc[1, 3])
day2_low_value = str(data2.iloc[2, 3])
day3_low_value = str(data2.iloc[3, 3])
day4_low_value = str(data2.iloc[4, 3])
day5_low_value = str(data2.iloc[5, 3])
day6_low_value = str(data2.iloc[6, 3])
day7_low_value = str(data2.iloc[6, 3]-1)

#WEATHER CONDITIONS
day1_weather = str(data2.iloc[0, 4]).lower()
day2_weather = str(data2.iloc[1, 4]).lower()
day3_weather = str(data2.iloc[2, 4]).lower()
day4_weather = str(data2.iloc[3, 4]).lower()
day5_weather = str(data2.iloc[4, 4]).lower()
day6_weather = str(data2.iloc[5, 4]).lower()
day7_weather = str(data2.iloc[6, 4]).lower()

slide = presentation.slides[slide_index3]

# Assign image file paths based on weather values using the dictionary mapping
day1_image_file = base_directory + weather_image_mapping.get(day1_weather, "Wind.png")
day2_image_file = base_directory + weather_image_mapping.get(day2_weather, "Wind.png")
day3_image_file = base_directory + weather_image_mapping.get(day3_weather, "Wind.png")
day4_image_file = base_directory + weather_image_mapping.get(day4_weather, "Wind.png")
day5_image_file = base_directory + weather_image_mapping.get(day5_weather, "Wind.png")
day6_image_file = base_directory + weather_image_mapping.get(day6_weather, "Wind.png")
day7_image_file = base_directory + weather_image_mapping.get(day7_weather, "Wind.png")



#SELECT AND MAKE TEXT BOX ACTIVE
day1_high = 9
textbox17 = slide.shapes[day1_high].text_frame
day2_high = 11
textbox18 = slide.shapes[day2_high].text_frame
day3_high = 10
textbox19 = slide.shapes[day3_high].text_frame
day4_high = 12
textbox20 = slide.shapes[day4_high].text_frame
day5_high = 14  
textbox21 = slide.shapes[day5_high].text_frame
day6_high = 13
textbox22 = slide.shapes[day6_high].text_frame
day7_high = 15  
textbox23 = slide.shapes[day7_high].text_frame

day1_low = 16
textbox24 = slide.shapes[day1_low].text_frame
day2_low = 18
textbox25 = slide.shapes[day2_low].text_frame
day3_low = 17
textbox26 = slide.shapes[day3_low].text_frame
day4_low = 19
textbox27 = slide.shapes[day4_low].text_frame
day5_low = 21  
textbox28 = slide.shapes[day5_low].text_frame
day6_low = 20
textbox29 = slide.shapes[day6_low].text_frame
day7_low = 22  
textbox30 = slide.shapes[day7_low].text_frame

#CLEAR TEXT, ADD NEW VALUES
textbox17.clear()
textbox17.text = day1_high_value
textbox18.clear()
textbox18.text = day2_high_value
textbox19.clear()
textbox19.text = day3_high_value
textbox20.clear()
textbox20.text = day4_high_value
textbox21.clear()
textbox21.text = day5_high_value
textbox22.clear()
textbox22.text = day6_high_value
textbox23.clear()
textbox23.text = day7_high_value

textbox24.clear()
textbox24.text = day1_low_value
textbox25.clear()
textbox25.text = day2_low_value
textbox26.clear()
textbox26.text = day3_low_value
textbox27.clear()
textbox27.text = day4_low_value
textbox28.clear()
textbox28.text = day5_low_value
textbox29.clear()
textbox29.text = day6_low_value
textbox30.clear()
textbox30.text = day7_low_value

# Iterate through the slides and replace images
for slide in presentation.slides:
    for shape in slide.shapes:
        if shape.name == "day1_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(day1_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "day2_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(day2_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "day3_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(day3_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "day4_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(day4_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "day5_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(day5_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "day6_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(day6_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "day7_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(day7_image_file, shape.left, shape.top, shape.width, shape.height)

#FORMATTING NEW TEXT


for paragraph in textbox17.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(50)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = False  
    paragraph.alignment = PP_ALIGN.CENTER


for paragraph in textbox18.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(50)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = False  
    paragraph.alignment = PP_ALIGN.CENTER


for paragraph in textbox19.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(50)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = False 
    paragraph.alignment = PP_ALIGN.CENTER


for paragraph in textbox20.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(50)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = False  
    paragraph.alignment = PP_ALIGN.CENTER


for paragraph in textbox21.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(50)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = False 
    paragraph.alignment = PP_ALIGN.CENTER


for paragraph in textbox22.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(50)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = False  
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox23.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(50)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = False  
    paragraph.alignment = PP_ALIGN.CENTER

# FORMATTING LOW TEMPERATURES


for paragraph in textbox24.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(38)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(180, 199, 231)
        run.font.bold = False  
    paragraph.alignment = PP_ALIGN.CENTER


for paragraph in textbox25.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(38)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(180, 199, 231)
        run.font.bold = False  
    paragraph.alignment = PP_ALIGN.CENTER


for paragraph in textbox26.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(38)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(180, 199, 231)
        run.font.bold = False 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox27.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(38)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(180, 199, 231)
        run.font.bold = False  
    paragraph.alignment = PP_ALIGN.CENTER


for paragraph in textbox28.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(38)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(180, 199, 231)
        run.font.bold = False 
    paragraph.alignment = PP_ALIGN.CENTER


for paragraph in textbox29.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(38)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(180, 199, 231)
        run.font.bold = False  
    paragraph.alignment = PP_ALIGN.CENTER


for paragraph in textbox30.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(38)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(180, 199, 231)
        run.font.bold = False  
    paragraph.alignment = PP_ALIGN.CENTER


slide = presentation.slides[slide_index4]

csv_file_path3 = "/Users/erick/Desktop/day_part_data.csv"
data3 = pd.read_csv(csv_file_path3)


#ASSIGN HIGH TEMP VALUES
daypart1_value = str(data3.iloc[1, 1])
daypart2_value = str(data2.iloc[0, 2])
daypart3_value = str(data3.iloc[1, 6])

#WEATHER CONDITIONS
daypart1_weather = str(data3.iloc[0,1])
daypart2_weather = str(data3.iloc[0,3])
daypart3_weather = str(data3.iloc[0,6])

#SELECT TEXT BOXES
daypart1_temp = 12
textbox101 = slide.shapes[daypart1_temp].text_frame
daypart2_temp = 13
textbox102 = slide.shapes[daypart2_temp].text_frame
daypart3_temp = 14
textbox103 = slide.shapes[daypart3_temp].text_frame

#CLEAR AND ADD NEW VALUES
textbox101.clear()
textbox101.text = daypart1_value
textbox102.clear()
textbox102.text = daypart2_value
textbox103.clear()
textbox103.text = daypart3_value


for paragraph in textbox101.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(66)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = False  
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox102.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(66)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = False  
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox103.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(66)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = False  
    paragraph.alignment = PP_ALIGN.CENTER

# Define the base directory for the image files
base_directory = "/Users/erick/Desktop/Weather_Graphics/Simple Weather Icons/weather_icons/"
base_directory2 = "/Users/erick/Desktop/Weather_Graphics/Simple Weather Icons/weather_icons_night"

weather_image_mapping2 = {
    "sky is clear": "Moon + Stars.png",
    "moderate rain": "Rain.png",
    "light rain": "Rain.png",
    "overcast clouds": "Cloud.png",
    "scattered clouds": "Night + Clouds.png",
    "broken clouds": "Night + Clouds.png",
    "few clouds": "Moon + Stars.png",
    "heavy intensity rain": "Thunderstorm 2.png",
    "clear sky": "Moon + Stars.png",
    "partly cloudy": "Night + Clouds.png",
    "sunny": "Moon + Stars.png",
    "patchy rain possible": "Rain.png",
    "heavy rain": "Thunderstorm 2.png",
    "thunderstorm with rain": "Thunderstorm 2.png",
    "thunderstorm with heavy rain": "Thunderstorm 2.png"

    # Add more mappings for other weather conditions
}

# Assign image file paths based on weather values using the dictionary mapping
daypart1_image_file = base_directory + weather_image_mapping.get(daypart1_weather, "Wind.png")
daypart2_image_file = base_directory + weather_image_mapping.get(daypart2_weather, "Wind.png")
daypart3_image_file = base_directory2 + weather_image_mapping2.get(daypart3_weather, "Wind.png")



# Iterate through the slides and replace images
for slide in presentation.slides:
    for shape in slide.shapes:
        if shape.name == "daypart1_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(daypart1_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "daypart2_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(daypart2_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "daypart3_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(daypart3_image_file, shape.left, shape.top, shape.width, shape.height)


# SLIDE NUMBER 10 BEGINS HERE **************************


slide = presentation.slides[slide_index5]


albany = 24  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox = slide.shapes[albany].text_frame

buffalo = 15  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox2 = slide.shapes[buffalo].text_frame

syracuse = 23  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox3 = slide.shapes[syracuse].text_frame

boston = 25  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox4 = slide.shapes[boston].text_frame

newyork = 13  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox5 = slide.shapes[newyork].text_frame

philadelphia = 21  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox6 = slide.shapes[philadelphia].text_frame

columbus = 17  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox7 = slide.shapes[columbus].text_frame

washington = 19  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox8 = slide.shapes[washington].text_frame


# In[253]:

#CLEAR TEXT, ADD NEW VALUES
textbox.clear()
textbox.text = f'{rain_value_alb}"'
textbox2.clear()
textbox2.text = f'{rain_value_buf}"'
textbox3.clear()
textbox3.text = f'{rain_value_syr}"'
textbox4.clear()
textbox4.text = f'{rain_value_bos}"'
textbox5.clear()
textbox5.text = f'{rain_value_nyc}"'
textbox6.clear()
textbox6.text = f'{rain_value_phl}"'
textbox7.clear()
textbox7.text = f'{rain_value_col}"'
textbox8.clear()
textbox8.text = f'{rain_value_was}"'

#FORMATTING NEW TEXT


for paragraph in textbox.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox2.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox3.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox4.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox5.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER


for paragraph in textbox6.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox7.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox8.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(0, 32, 96)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER



# SLIDE NUMBER 18 BEGINS HERE **************************


slide = presentation.slides[slide_index6]


albany2 = 24  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox = slide.shapes[albany2].text_frame

buffalo2 = 11  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox2 = slide.shapes[buffalo2].text_frame

syracuse2 = 23  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox3 = slide.shapes[syracuse2].text_frame

boston2 = 25  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox4 = slide.shapes[boston2].text_frame

newyork2 = 8  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox5 = slide.shapes[newyork2].text_frame

philadelphia2 = 20  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox6 = slide.shapes[philadelphia2].text_frame

columbus2 = 14  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox7 = slide.shapes[columbus2].text_frame

washington2 = 17  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox8 = slide.shapes[washington2].text_frame


# In[253]:

#CLEAR TEXT, ADD NEW VALUES
textbox.clear()
textbox.text = app_value_alb
textbox2.clear()
textbox2.text = app_value_buf
textbox3.clear()
textbox3.text = app_value_syr
textbox4.clear()
textbox4.text = app_value_bos
textbox5.clear()
textbox5.text = app_value_nyc
textbox6.clear()
textbox6.text = app_value_phl
textbox7.clear()
textbox7.text = app_value_col
textbox8.clear()
textbox8.text = app_value_was

#FORMATTING NEW TEXT


for paragraph in textbox.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(139, 0, 0)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox2.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(139, 0, 0)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox3.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(139, 0, 0)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox4.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(139, 0, 0)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox5.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(139, 0, 0)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER


for paragraph in textbox6.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(139, 0, 0)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox7.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(139, 0, 0)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox8.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(48)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(139, 0, 0)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

#IMAGE CODING HERE


# Define the base directory for the image files
base_directory = "/Users/erick/Desktop/Weather_Graphics/Simple Weather Icons/weather_icons/"

# Define the dictionary mapping weather values to image file paths
weather_image_mapping = {
    "sky is clear": "Sun 3.png",
    "moderate rain": "Rain.png",
    "light rain": "Rain + Sun.png",
    "overcast clouds": "Cloud.png",
    "scattered clouds": "Sun & Clouds.png",
    "broken clouds": "Sun & Clouds.png",
    "few clouds": "Sun 3.png",
    "heavy intensity rain": "Thunderstorm & Sun.png",
    "clear sky": "Sun 3.png",
    "partly cloudy": "Sun & Clouds.png",
    "sunny": "Sun 3.png",
    "patchy rain possible": "Rain + Sun.png",
    "heavy rain": "Thunderstorm & Sun.png",
    "thunderstorm with rain": "Thunderstorm & Sun.png",
    "thunderstorm with heavy rain": "Thunderstorm 2.png"

    # Add more mappings for other weather conditions
}

# Assign image file paths based on weather values using the dictionary mapping
nyc_image_file = base_directory + weather_image_mapping.get(nyc_weather, "Wind.png")
alb_image_file = base_directory + weather_image_mapping.get(alb_weather, "Wind.png")
bos_image_file = base_directory + weather_image_mapping.get(bos_weather, "Wind.png")
buf_image_file = base_directory + weather_image_mapping.get(buf_weather, "Wind.png")
syr_image_file = base_directory + weather_image_mapping.get(syr_weather, "Wind.png")
phl_image_file = base_directory + weather_image_mapping.get(phl_weather, "Wind.png")
was_image_file = base_directory + weather_image_mapping.get(was_weather, "Wind.png")
col_image_file = base_directory + weather_image_mapping.get(col_weather, "Wind.png")


# Iterate through the slides and replace images
for slide in presentation.slides:
    for shape in slide.shapes:
        if shape.name == "nyc_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(nyc_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "alb_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(alb_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "bos_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(bos_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "buf_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(buf_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "syr_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(syr_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "phl_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(phl_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "was_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(was_image_file, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "col_icon":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(col_image_file, shape.left, shape.top, shape.width, shape.height)


#WEEKEND WEATHER SLIDE 17 STARTS HERE
slide = presentation.slides[slide_index7]

albany_sat = 9  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox = slide.shapes[albany_sat].text_frame
albany_sun = 10  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox2 = slide.shapes[albany_sun].text_frame

buffalo_sat = 45  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox3 = slide.shapes[buffalo_sat].text_frame
buffalo_sat = 46  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox4 = slide.shapes[buffalo_sat].text_frame

boston_sat = 18  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox5 = slide.shapes[boston_sat].text_frame
boston_sun = 19  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox6 = slide.shapes[boston_sun].text_frame

newyork_sat = 27  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox7 = slide.shapes[newyork_sat].text_frame
newyork_sun = 28  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox8 = slide.shapes[newyork_sun].text_frame

philadelphia_sat = 36  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox9 = slide.shapes[philadelphia_sat].text_frame
philadelphia_sun = 37  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox10 = slide.shapes[philadelphia_sun].text_frame

columbus_sat = 54  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox11 = slide.shapes[columbus_sat].text_frame
columbus_sun = 55  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox12 = slide.shapes[columbus_sun].text_frame

washington_sat = 63  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox13 = slide.shapes[washington_sat].text_frame
washington_sun = 64  # Textbox index is 0-based, so textbox 9 corresponds to index 8
textbox14 = slide.shapes[washington_sun].text_frame

textbox.clear()
textbox.text = alb_sat_temp
textbox2.clear()
textbox2.text = alb_sun_temp
textbox3.clear()
textbox3.text = buf_sat_temp
textbox4.clear()
textbox4.text = buf_sun_temp
textbox5.clear()
textbox5.text = bos_sat_temp
textbox6.clear()
textbox6.text = bos_sun_temp
textbox7.clear()
textbox7.text = nyc_sat_temp
textbox8.clear()
textbox8.text = nyc_sun_temp
textbox9.clear()
textbox9.text = phl_sat_temp
textbox10.clear()
textbox10.text = phl_sun_temp
textbox11.clear()
textbox11.text = col_sat_temp
textbox12.clear()
textbox12.text = col_sun_temp
textbox13.clear()
textbox13.text = was_sat_temp
textbox14.clear()
textbox14.text = was_sun_temp

for paragraph in textbox.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(22)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox2.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(22)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox3.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(22)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox4.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(22)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox5.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(22)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox6.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(22)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox7.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(22)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox8.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(22)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox9.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(22)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox10.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(22)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox11.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(22)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox12.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(22)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER

for paragraph in textbox13.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(22)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER


for paragraph in textbox14.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(22)  # Set font size to 48 points
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.bold = True 
    paragraph.alignment = PP_ALIGN.CENTER


# Assign image file paths based on weather values using the dictionary mapping
nyc_image_file = base_directory + weather_image_mapping.get(nyc_weather, "Wind.png")
alb_image_file = base_directory + weather_image_mapping.get(alb_weather, "Wind.png")
bos_image_file = base_directory + weather_image_mapping.get(bos_weather, "Wind.png")
buf_image_file = base_directory + weather_image_mapping.get(buf_weather, "Wind.png")
syr_image_file = base_directory + weather_image_mapping.get(syr_weather, "Wind.png")
phl_image_file = base_directory + weather_image_mapping.get(phl_weather, "Wind.png")
was_image_file = base_directory + weather_image_mapping.get(was_weather, "Wind.png")
col_image_file = base_directory + weather_image_mapping.get(col_weather, "Wind.png")


# Assign image file paths based on weather values using the dictionary mapping
nyc_sat = base_directory + weather_image_mapping.get(nyc_sat, "Wind.png")
alb_sat = base_directory + weather_image_mapping.get(alb_sat, "Wind.png")
bos_sat = base_directory + weather_image_mapping.get(bos_sat, "Wind.png")
buf_sat = base_directory + weather_image_mapping.get(buf_sat, "Wind.png")
phl_sat = base_directory + weather_image_mapping.get(phl_sat, "Wind.png")
was_sat = base_directory + weather_image_mapping.get(was_sat, "Wind.png")
col_sat = base_directory + weather_image_mapping.get(col_sat, "Wind.png")
nyc_sun = base_directory + weather_image_mapping.get(nyc_sun, "Wind.png")
alb_sun = base_directory + weather_image_mapping.get(alb_sun, "Wind.png")
bos_sun = base_directory + weather_image_mapping.get(bos_sun, "Wind.png")
buf_sun = base_directory + weather_image_mapping.get(buf_sun, "Wind.png")
phl_sun = base_directory + weather_image_mapping.get(phl_sun, "Wind.png")
was_sun = base_directory + weather_image_mapping.get(was_sun, "Wind.png")
col_sun = base_directory + weather_image_mapping.get(col_sun, "Wind.png")


for slide in presentation.slides:
    for shape in slide.shapes:
        if shape.name == "washington_sat":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(was_sat, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "albany_sat":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(alb_sat, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "boston_sat":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(bos_sat, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "buffalo_sat":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(buf_sat, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "philadelphia_sat":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(phl_sat, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "washington_sat":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(was_sat, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "columbus_sat":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(col_sat, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "albany_sun":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(alb_sun, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "boston_sun":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(bos_sun, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "buffalo_sun":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(buf_sun, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "philadelphia_sun":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(phl_sun, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "washington_sun":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(was_sun, shape.left, shape.top, shape.width, shape.height)
        elif shape.name == "columbus_sun":
            shape.element.getparent().remove(shape.element)
            slide.shapes.add_picture(col_sun, shape.left, shape.top, shape.width, shape.height)

#THIS UPDATES THE PRESENTATION
updated_powerpoint_file_path = "/Users/erick/Desktop/Weather_Update.pptx"
presentation.save(updated_powerpoint_file_path)

